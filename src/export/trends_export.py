"""Trends-specific export for VPG Intelligence Digest V2.4.

Generates:
- Excel workbook with tabs: Trend Alerts, Industry Momentum,
  Signal Volume by Week, Competitor Trends, Keyword Momentum
- PowerPoint with slides: Title, What's Moving, Industry Momentum,
  Signal Volume (table), Competitor Trends
"""

import logging
from datetime import datetime
from io import BytesIO

from src.db import get_connection
from src.trends.tracker import (
    get_trend_alerts,
    get_industry_momentum,
    get_signal_volume_over_time,
    get_competitor_trends,
    get_trend_summary,
)

logger = logging.getLogger(__name__)

# VPG brand colors
VPG_NAVY = (27, 42, 74)
VPG_BLUE = (46, 117, 182)
VPG_ACCENT = (232, 121, 47)
WHITE = (255, 255, 255)


def export_trends_excel(bu_code: str | None = None,
                        start_date: str | None = None,
                        end_date: str | None = None) -> BytesIO:
    """Export Trends data as Excel workbook."""
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    wb = openpyxl.Workbook()
    navy_fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
    white_font = Font(name="Arial", bold=True, color="FFFFFF", size=10)
    center = Alignment(horizontal="center", vertical="center")

    def style_header(ws):
        for cell in ws[1]:
            cell.fill = navy_fill
            cell.font = white_font
            cell.alignment = center

    def auto_width(ws):
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=0)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

    conn = get_connection()
    try:
        # Tab 1: Trend Alerts
        ws1 = wb.active
        ws1.title = "Trend Alerts"
        ws1.append(["Trend Name", "Type", "Change %", "Signals", "Period (weeks)",
                     "Companies", "Top Signal", "Suggested Action", "Industry"])
        alerts = get_trend_alerts(conn=conn, bu_code=bu_code, limit=20)
        for a in alerts:
            companies = ", ".join(a.get("companies", []) or [])
            ws1.append([
                a.get("trend_name", ""),
                a.get("trend_type", ""),
                a.get("change_percent", 0),
                a.get("signal_count", 0),
                a.get("period_weeks", 0),
                companies,
                a.get("top_signal_headline", ""),
                a.get("suggested_action", ""),
                a.get("industry", ""),
            ])
        style_header(ws1)
        auto_width(ws1)

        # Tab 2: Industry Momentum
        ws2 = wb.create_sheet("Industry Momentum")
        ws2.append(["Industry", "Signals", "Avg Score", "Score Change",
                     "Opportunities", "Threats", "Sentiment", "Change %", "Top Competitor"])
        industries = get_industry_momentum(conn=conn, bu_code=bu_code,
                                            start_date=start_date, end_date=end_date)
        for ind in industries:
            ws2.append([
                ind["name"], ind["signal_count"], ind["avg_score"],
                ind.get("score_change", 0), ind["opportunities"], ind["threats"],
                ind["sentiment"], ind.get("change_percent", 0),
                ind.get("top_competitor", ""),
            ])
        style_header(ws2)
        auto_width(ws2)

        # Tab 3: Signal Volume by Week
        ws3 = wb.create_sheet("Signal Volume by Week")
        volume = get_signal_volume_over_time(conn=conn, bu_code=bu_code)
        week_labels = volume.get("weeks", [])
        series = volume.get("series", [])
        ws3.append(["Week"] + [s["name"] for s in series])
        for i, wl in enumerate(week_labels):
            row = [wl]
            for s in series:
                dp = s["data"][i] if i < len(s["data"]) else {"count": 0}
                row.append(dp.get("count", 0))
            ws3.append(row)
        style_header(ws3)
        auto_width(ws3)

        # Tab 4: Competitor Trends
        ws4 = wb.create_sheet("Competitor Trends")
        ws4.append(["Competitor", "This Period", "Prior Period", "Change %", "Trend"])
        competitors = get_competitor_trends(conn=conn, bu_code=bu_code,
                                             start_date=start_date, end_date=end_date)
        for c in competitors:
            ws4.append([
                c["name"], c["this_period"], c["prior_period"],
                c["change_percent"], c["trend"],
            ])
        style_header(ws4)
        auto_width(ws4)

        # Tab 5: Keyword Momentum
        ws5 = wb.create_sheet("Keyword Momentum")
        ws5.append(["Trend", "Type", "Momentum", "Signals", "WoW %",
                     "Avg Score", "Max Score", "First Seen", "Last Seen"])
        summary = get_trend_summary(conn=conn, limit=100)
        for t in summary.get("trends", []):
            ws5.append([
                t["label"], t["type"], t["momentum"], t["count"],
                t.get("change_pct", 0), t.get("avg_score", 0),
                t.get("max_score", 0), t.get("first_seen", ""),
                t.get("last_seen", ""),
            ])
        style_header(ws5)
        auto_width(ws5)

    finally:
        conn.close()

    buffer = BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer


def export_trends_pptx(bu_code: str | None = None,
                       start_date: str | None = None,
                       end_date: str | None = None) -> BytesIO:
    """Export Trends data as PowerPoint presentation."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    conn = get_connection()
    try:
        bu_label = bu_code or "All BUs"
        date_label = ""
        if start_date and end_date:
            date_label = f" — {start_date} to {end_date}"
        elif start_date:
            date_label = f" — From {start_date}"

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*VPG_NAVY)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Trends & Strategic Intelligence"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*WHITE)
        p2 = tf.add_paragraph()
        p2.text = f"{bu_label}{date_label} — {datetime.now().strftime('%B %d, %Y')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(*VPG_BLUE)

        footer = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
        fp = footer.text_frame.paragraphs[0]
        fp.text = "Vishay Precision Group | Confidential"
        fp.font.size = Pt(10)
        fp.font.color.rgb = RGBColor(150, 150, 180)

        # Slide 2: What's Moving (Trend Alerts)
        alerts = get_trend_alerts(conn=conn, bu_code=bu_code, limit=5)
        if alerts:
            slide2 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = "What's Moving"
            tp.font.size = Pt(24)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(*VPG_NAVY)

            y = 1.2
            for alert in alerts[:5]:
                box = slide2.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(1.0))
                tf = box.text_frame
                tf.word_wrap = True

                type_icons = {"rising": "RISING", "declining": "DECLINING",
                              "new": "NEW", "persistent": "PERSISTENT"}
                icon = type_icons.get(alert.get("trend_type", ""), "")

                p = tf.paragraphs[0]
                run1 = p.add_run()
                run1.text = f"{icon}: {alert.get('trend_name', '')}  "
                run1.font.size = Pt(12)
                run1.font.bold = True
                run1.font.color.rgb = RGBColor(*VPG_ACCENT)

                change = alert.get("change_percent", 0)
                arrow = "↑" if change > 0 else "↓" if change < 0 else "→"
                run2 = p.add_run()
                run2.text = f"{arrow}{abs(change):.0f}%"
                run2.font.size = Pt(12)
                run2.font.color.rgb = RGBColor(*VPG_BLUE)

                p2 = tf.add_paragraph()
                companies = ", ".join(alert.get("companies", [])[:3])
                p2.text = (f"{alert.get('signal_count', 0)} signals over "
                           f"{alert.get('period_weeks', 0)} weeks"
                           f"{f' — Key: {companies}' if companies else ''}")
                p2.font.size = Pt(10)
                p2.font.color.rgb = RGBColor(80, 80, 80)

                y += 1.15

        # Slide 3: Industry Momentum
        industries = get_industry_momentum(conn=conn, bu_code=bu_code,
                                            start_date=start_date, end_date=end_date)
        if industries:
            slide3 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = "Industry Momentum"
            tp.font.size = Pt(24)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(*VPG_NAVY)

            # Table
            display = industries[:10]
            rows = len(display) + 1
            cols = 6
            tbl = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2),
                                           Inches(9), Inches(0.35 * rows))
            table = tbl.table
            for i, h in enumerate(["Industry", "Signals", "Avg Score", "Sentiment",
                                    "Change %", "Top Competitor"]):
                cell = table.cell(0, i)
                cell.text = h
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(9)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(*WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*VPG_NAVY)

            for r, ind in enumerate(display, 1):
                vals = [ind["name"], str(ind["signal_count"]),
                        f"{ind['avg_score']:.1f}", ind["sentiment"],
                        f"{ind.get('change_percent', 0):+.0f}%",
                        ind.get("top_competitor", "—")]
                for c, val in enumerate(vals):
                    table.cell(r, c).text = val
                    for para in table.cell(r, c).text_frame.paragraphs:
                        para.font.size = Pt(9)

        # Slide 4: Signal Volume (table form since we can't render charts in pptx easily)
        volume = get_signal_volume_over_time(conn=conn, bu_code=bu_code)
        if volume.get("series"):
            slide4 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = "Signal Volume Over Time"
            tp.font.size = Pt(24)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(*VPG_NAVY)

            weeks = volume["weeks"]
            series = volume["series"][:6]  # Limit for slide space
            rows = len(series) + 1
            cols = min(len(weeks) + 1, 14)  # Limit columns
            display_weeks = weeks[:cols - 1]

            tbl = slide4.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2),
                                           Inches(9.4), Inches(0.35 * rows))
            table = tbl.table
            table.cell(0, 0).text = "Series"
            for i, w in enumerate(display_weeks):
                table.cell(0, i + 1).text = w
            for i in range(cols):
                cell = table.cell(0, i)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(8)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(*WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*VPG_NAVY)

            for r, s in enumerate(series, 1):
                table.cell(r, 0).text = s["name"]
                for c, dp in enumerate(s["data"][:len(display_weeks)]):
                    table.cell(r, c + 1).text = str(dp.get("count", 0))
                for c in range(cols):
                    for para in table.cell(r, c).text_frame.paragraphs:
                        para.font.size = Pt(8)

        # Slide 5: Competitor Trends
        competitors = get_competitor_trends(conn=conn, bu_code=bu_code,
                                             start_date=start_date, end_date=end_date)
        if competitors:
            slide5 = prs.slides.add_slide(prs.slide_layouts[6])
            title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            tp = title_box.text_frame.paragraphs[0]
            tp.text = "Competitor Trends"
            tp.font.size = Pt(24)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(*VPG_NAVY)

            display = competitors[:12]
            rows = len(display) + 1
            cols = 5
            tbl = slide5.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2),
                                           Inches(9), Inches(0.35 * rows))
            table = tbl.table
            for i, h in enumerate(["Competitor", "This Period", "Prior Period",
                                    "Change %", "Trend"]):
                cell = table.cell(0, i)
                cell.text = h
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(9)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(*WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*VPG_NAVY)

            trend_icons = {"rising": "📈", "declining": "📉", "stable": "➡️"}
            for r, c in enumerate(display, 1):
                vals = [c["name"], str(c["this_period"]), str(c["prior_period"]),
                        f"{c['change_percent']:+.0f}%",
                        trend_icons.get(c["trend"], "➡️")]
                for ci, val in enumerate(vals):
                    table.cell(r, ci).text = val
                    for para in table.cell(r, ci).text_frame.paragraphs:
                        para.font.size = Pt(9)

    finally:
        conn.close()

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
