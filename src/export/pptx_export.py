"""PowerPoint export for VPG Intelligence Digest.

Generates .pptx presentations with:
- Title slide with VPG branding
- Executive summary slide
- Top signals slides (one per high-scoring signal)
- BU breakdown slide
- Trends overview slide

Uses python-pptx for generation.
Falls back gracefully if python-pptx is not installed.
"""

import logging
from datetime import datetime
from io import BytesIO
from pathlib import Path

from src.config import get_business_units
from src.db import get_connection

logger = logging.getLogger(__name__)

# VPG brand colors (RGB tuples)
VPG_NAVY = (27, 42, 74)
VPG_BLUE = (46, 117, 182)
VPG_ACCENT = (232, 121, 47)
WHITE = (255, 255, 255)
LIGHT_GRAY = (240, 240, 240)

# Signal type display config
SIGNAL_TYPE_LABELS = {
    "competitive-threat": "Competitive Threat",
    "revenue-opportunity": "Revenue Opportunity",
    "market-shift": "Market Shift",
    "partnership-signal": "Partnership Signal",
    "customer-intelligence": "Customer Intelligence",
    "technology-trend": "Technology Trend",
    "trade-tariff": "Trade & Tariff",
}


def _get_pptx():
    """Import python-pptx or return None."""
    try:
        import pptx
        return pptx
    except ImportError:
        logger.warning("python-pptx not installed — PowerPoint export disabled. pip install python-pptx")
        return None


def _rgb(r, g, b):
    """Create an RGBColor from tuple."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _add_title_slide(prs, pptx_mod):
    """Add a VPG-branded title slide."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*VPG_NAVY)

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "VPG Intelligence Digest"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*WHITE)

    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = f"Weekly Intelligence Report — {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(*VPG_BLUE)

    # Footer
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "Vishay Precision Group | Confidential"
    p3.font.size = Pt(10)
    p3.font.color.rgb = RGBColor(150, 150, 180)

    return slide


def _add_summary_slide(prs, pptx_mod, stats):
    """Add an executive summary slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*VPG_NAVY)

    # Stats grid
    y = 1.3
    stat_items = [
        ("Signals Analyzed", str(stats.get("total_signals", 0))),
        ("High-Impact Signals (8.0+)", str(stats.get("high_impact", 0))),
        ("Business Units with Signals", str(stats.get("bus_with_signals", 0))),
        ("Rising Trends", str(stats.get("rising_trends", 0))),
        ("Top Signal Type", stats.get("top_signal_type", "N/A")),
        ("Average Score", f"{stats.get('avg_score', 0):.1f}"),
    ]

    for label, value in stat_items:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(5), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{label}: "
        run1.font.size = Pt(14)
        run1.font.color.rgb = RGBColor(100, 100, 100)
        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(14)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor(*VPG_NAVY)
        y += 0.55

    return slide


def _add_signal_slide(prs, pptx_mod, signal):
    """Add a slide for a single high-impact signal."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Signal type badge + score
    type_label = SIGNAL_TYPE_LABELS.get(signal.get("signal_type", ""), signal.get("signal_type", ""))
    score = signal.get("score_composite", 0)

    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7), Inches(0.5))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"{type_label}  "
    run1.font.size = Pt(11)
    run1.font.color.rgb = RGBColor(*VPG_ACCENT)
    run1.font.bold = True
    run2 = p.add_run()
    run2.text = f"Score: {score:.1f}"
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*VPG_BLUE)

    # Headline
    headline_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    tf2 = headline_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = signal.get("headline", signal.get("title", ""))
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(*VPG_NAVY)

    # Content sections
    y = 1.8
    sections = [
        ("WHAT", signal.get("what_summary", "")),
        ("WHY IT MATTERS", signal.get("why_it_matters", "")),
        ("QUICK WIN", signal.get("quick_win", "")),
        ("OWNER", signal.get("suggested_owner", "")),
        ("EST. IMPACT", signal.get("estimated_impact", "")),
    ]

    for label, content in sections:
        if not content:
            continue
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{label}: "
        run1.font.size = Pt(11)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(*VPG_BLUE)
        run2 = p.add_run()
        run2.text = str(content)[:300]
        run2.font.size = Pt(11)
        run2.font.color.rgb = RGBColor(60, 60, 60)
        y += 0.7

    # Source links
    source_links = signal.get("source_links", [])
    if source_links:
        links_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
        tf = links_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "SOURCES: "
        run1.font.size = Pt(9)
        run1.font.bold = True
        run1.font.color.rgb = RGBColor(100, 100, 100)
        run2 = p.add_run()
        run2.text = " | ".join(source_links[:5])
        run2.font.size = Pt(8)
        run2.font.color.rgb = RGBColor(46, 117, 182)

    # BUs
    bu_str = ", ".join(signal.get("bus", []))
    if bu_str:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.3))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Business Units: {bu_str}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(130, 130, 130)

    return slide


def _add_bu_breakdown_slide(prs, pptx_mod, bu_data):
    """Add a BU breakdown summary slide."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Unit Signal Breakdown"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*VPG_NAVY)

    # Table
    rows = len(bu_data) + 1
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * rows))
    table = table_shape.table

    # Headers
    for i, header in enumerate(["Business Unit", "Signals", "Avg Score", "Top Type"]):
        cell = table.cell(0, i)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(*WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*VPG_NAVY)

    # Data rows
    for r, bu in enumerate(bu_data, 1):
        table.cell(r, 0).text = bu["name"]
        table.cell(r, 1).text = str(bu["signal_count"])
        table.cell(r, 2).text = f"{bu['avg_score']:.1f}"
        table.cell(r, 3).text = bu.get("top_type", "N/A")
        for c in range(cols):
            for paragraph in table.cell(r, c).text_frame.paragraphs:
                paragraph.font.size = Pt(10)
            if r % 2 == 0:
                table.cell(r, c).fill.solid()
                table.cell(r, c).fill.fore_color.rgb = RGBColor(*LIGHT_GRAY)

    return slide


def export_signals_pptx(start_date: str | None = None, end_date: str | None = None,
                        bu_id: str | None = None, signal_type: str | None = None,
                        industry_id: str | None = None, min_score: float = 0,
                        max_signals: int = 10,
                        output_path: Path | None = None) -> BytesIO | Path:
    """Export intelligence to a PowerPoint presentation.

    Args:
        start_date: Optional YYYY-MM-DD start filter.
        end_date: Optional YYYY-MM-DD end filter.
        max_signals: Maximum signal slides to include.
        output_path: If provided, saves to file and returns Path.

    Returns:
        BytesIO buffer or Path to saved file.
    """
    pptx_mod = _get_pptx()
    if pptx_mod is None:
        raise ImportError("python-pptx is required for PowerPoint export")

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    conn = get_connection()
    try:
        # Fetch signals
        query = """
            SELECT s.id, s.title, sa.signal_type, sa.headline, sa.score_composite,
                   sa.what_summary, sa.why_it_matters, sa.quick_win,
                   sa.suggested_owner, sa.estimated_impact, sa.validation_level,
                   s.url
            FROM signals s
            JOIN signal_analysis sa ON s.id = sa.signal_id
            WHERE s.status IN ('scored', 'published')
              AND COALESCE(s.dismissed, 0) = 0
        """
        params = []
        if start_date:
            query += " AND s.collected_at >= ?"
            params.append(start_date)
        if end_date:
            query += " AND s.collected_at <= ?"
            params.append(end_date + " 23:59:59")
        if signal_type:
            query += " AND sa.signal_type = ?"
            params.append(signal_type)
        if min_score > 0:
            query += " AND sa.score_composite >= ?"
            params.append(min_score)
        if bu_id:
            query += " AND s.id IN (SELECT signal_id FROM signal_bus WHERE bu_id = ?)"
            params.append(bu_id)
        if industry_id:
            query += " AND s.id IN (SELECT signal_id FROM signal_industries WHERE industry_id = ?)"
            params.append(industry_id)
        query += " ORDER BY sa.score_composite DESC"

        rows = conn.execute(query, params).fetchall()

        # Enrich with BU data and source links
        signals = []
        for row in rows:
            bus = conn.execute(
                "SELECT bu_id FROM signal_bus WHERE signal_id = ?", (row[0],)
            ).fetchall()
            # Collect all source links
            source_links = [row[11]] if row[11] else []
            validations = conn.execute(
                "SELECT corroborating_url FROM signal_validations WHERE signal_id = ?",
                (row[0],)
            ).fetchall()
            source_links.extend(v[0] for v in validations if v[0])

            signals.append({
                "id": row[0], "title": row[1], "signal_type": row[2],
                "headline": row[3], "score_composite": row[4] or 0,
                "what_summary": row[5], "why_it_matters": row[6],
                "quick_win": row[7], "suggested_owner": row[8],
                "estimated_impact": row[9], "validation_level": row[10],
                "url": row[11],
                "bus": [b[0] for b in bus],
                "source_links": source_links,
            })

        # Compute stats
        from collections import Counter
        type_counter = Counter(s["signal_type"] for s in signals)
        bu_counter = Counter()
        bu_scores = {}
        bu_types = {}
        for s in signals:
            for bu in s["bus"]:
                bu_counter[bu] += 1
                bu_scores.setdefault(bu, []).append(s["score_composite"])
                bu_types.setdefault(bu, Counter())[s["signal_type"]] += 1

        bu_config = get_business_units()
        bu_names = {bu["id"]: bu["name"] for bu in bu_config.get("business_units", [])}

        stats = {
            "total_signals": len(signals),
            "high_impact": sum(1 for s in signals if s["score_composite"] >= 8.0),
            "bus_with_signals": len(bu_counter),
            "rising_trends": conn.execute(
                "SELECT COUNT(*) FROM trends WHERE momentum IN ('rising', 'spike')"
            ).fetchone()[0],
            "top_signal_type": SIGNAL_TYPE_LABELS.get(
                type_counter.most_common(1)[0][0], "N/A"
            ) if type_counter else "N/A",
            "avg_score": sum(s["score_composite"] for s in signals) / max(len(signals), 1),
        }

        bu_data = []
        for bu_id, count in bu_counter.most_common():
            scores = bu_scores.get(bu_id, [])
            top_type = bu_types.get(bu_id, Counter()).most_common(1)
            bu_data.append({
                "name": bu_names.get(bu_id, bu_id),
                "signal_count": count,
                "avg_score": sum(scores) / max(len(scores), 1),
                "top_type": SIGNAL_TYPE_LABELS.get(top_type[0][0], "N/A") if top_type else "N/A",
            })

        # Build slides
        _add_title_slide(prs, pptx_mod)
        _add_summary_slide(prs, pptx_mod, stats)

        # Top signal slides
        for signal in signals[:max_signals]:
            _add_signal_slide(prs, pptx_mod, signal)

        # BU breakdown
        if bu_data:
            _add_bu_breakdown_slide(prs, pptx_mod, bu_data)

    finally:
        conn.close()

    # Output
    if output_path:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PowerPoint export saved to %s", output_path)
        return output_path
    else:
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
