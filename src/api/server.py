"""FastAPI server for VPG Intelligence Digest management UI.

Provides REST API endpoints for:
- Recipient management (add, edit, remove)
- Business unit & industry configuration
- Source management (add, edit, toggle, remove)
- Pipeline execution (run digest, dry-run, pause/resume/cancel)
- Digest history and status
- Delivery schedule editing
- Trend data and history

Run with: python -m src.api.server
"""

import logging
import threading
from datetime import datetime
from io import BytesIO
from pathlib import Path

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from src.config import (
    CONFIG_DIR,
    MOCK_OUTPUT_DIR,
    PROJECT_ROOT,
    get_business_units,
    get_industries,
    get_recipients,
    get_scoring_weights,
    get_sources,
    save_business_units,
    save_industries,
    save_recipients,
    save_scoring_weights,
    save_sources,
)
from src.db import (
    get_connection,
    init_db,
    get_all_industries,
    upsert_industry,
    delete_industry as db_delete_industry,
    get_all_keywords,
    upsert_keyword,
    delete_keyword as db_delete_keyword,
    bulk_import_keywords,
    get_pipeline_runs_by_timeframe,
    get_signals_by_timeframe,
)

logger = logging.getLogger(__name__)

app = FastAPI(
    title="VPG Intelligence Digest",
    description="Management UI for the VPG Weekly Intelligence Digest",
    version="3.0.0",
)

# CORS — allow the React dev server and local access
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:5173", "http://127.0.0.1:3000", "http://127.0.0.1:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Track background pipeline runs
_pipeline_status = {
    "running": False,
    "paused": False,
    "last_result": None,
    "last_run": None,
    "current_stage": "",
}


@app.on_event("startup")
def _startup():
    """Initialize DB, seed industries/keywords from config, clean stale runs."""
    try:
        print("[STARTUP] Initializing database...")
        conn = get_connection()
        init_db()

        # Clean up stale pipeline runs (from prior crash/restart)
        conn.execute(
            "UPDATE pipeline_runs SET status = 'failed', completed_at = datetime('now') "
            "WHERE status = 'running'"
        )
        conn.commit()

        # Migrate FIRST: add dismissed/handled columns + reddit table (V2.3)
        print("[STARTUP] Running V2.3 migration...")
        _migrate_v23(conn)

        # Sync industries + keywords from config
        print("[STARTUP] Syncing industries from config...")
        _auto_seed_industries(conn)

        ind_count = conn.execute("SELECT COUNT(*) FROM industries").fetchone()[0]
        kw_count = conn.execute("SELECT COUNT(*) FROM keywords").fetchone()[0]
        sub_count = conn.execute("SELECT COUNT(*) FROM reddit_subreddits").fetchone()[0]
        print(f"[STARTUP] DB ready: {ind_count} industries, {kw_count} keywords, {sub_count} subreddits")

        conn.close()
    except Exception as e:
        print(f"[STARTUP] ERROR: {e}")
        import traceback
        traceback.print_exc()
        logger.warning("Startup initialization error: %s", e)


def _migrate_v23(conn):
    """Add V2.3 columns (dismissed, handled) to signals + reddit_subreddits table."""
    try:
        cols = {row[1] for row in conn.execute("PRAGMA table_info(signals)").fetchall()}
        for col, ddl in [
            ("dismissed", "INTEGER NOT NULL DEFAULT 0"),
            ("dismissed_at", "DATETIME"),
            ("handled", "INTEGER NOT NULL DEFAULT 0"),
            ("handled_at", "DATETIME"),
            ("handled_by", "TEXT DEFAULT ''"),
        ]:
            if col not in cols:
                conn.execute(f"ALTER TABLE signals ADD COLUMN {col} {ddl}")
        conn.execute("""
            CREATE TABLE IF NOT EXISTS reddit_subreddits (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                name TEXT NOT NULL UNIQUE,
                category TEXT DEFAULT '',
                active INTEGER NOT NULL DEFAULT 1,
                notes TEXT DEFAULT '',
                created_at DATETIME NOT NULL DEFAULT (datetime('now')),
                updated_at DATETIME NOT NULL DEFAULT (datetime('now'))
            )
        """)
        conn.commit()
        # Sync subreddits from DEFAULT_SUBREDDITS — insert any missing ones
        from src.collector.reddit_collector import DEFAULT_SUBREDDITS
        _CATEGORIES = {
            "robotics": "Robotics & Automation", "Automate": "Robotics & Automation", "ROS": "Robotics & Automation",
            "aerospace": "Aerospace & Defense", "DefenseIndustry": "Aerospace & Defense",
            "electricvehicles": "Automotive & EV", "SelfDrivingCars": "Automotive & EV", "automotive": "Automotive & EV",
            "metalworking": "Steel & Metals", "manufacturing": "Manufacturing", "Machinists": "Manufacturing",
            "electronics": "Sensors & Instrumentation", "sensors": "Sensors & Instrumentation", "ECE": "Sensors & Instrumentation",
            "materials": "Materials Science", "MaterialsScience": "Materials Science",
            "engineering": "Test & Measurement",
            "SupplyChain": "Trade & Tariffs", "Economics": "Trade & Tariffs",
            "mining": "Mining & Heavy Equipment", "HeavyEquipment": "Mining & Heavy Equipment",
            "agriculture": "Agriculture", "PrecisionAg": "Agriculture", "farming": "Agriculture",
            "SportsScience": "Sports & Performance", "Biomechanics": "Sports & Performance", "sportsmedicine": "Sports & Performance",
            "telecom": "Communication & Telecom", "5G": "Communication & Telecom", "rfelectronics": "Communication & Telecom",
            "Construction": "Infrastructure & Construction", "civilengineering": "Infrastructure & Construction", "infrastructure": "Infrastructure & Construction",
        }
        for sub in DEFAULT_SUBREDDITS:
            cat = _CATEGORIES.get(sub, "General")
            conn.execute(
                "INSERT OR IGNORE INTO reddit_subreddits (name, category) VALUES (?, ?)",
                (sub, cat)
            )
        conn.commit()
    except Exception as e:
        print(f"[STARTUP] V2.3 migration error: {e}")
        import traceback
        traceback.print_exc()
        logger.warning("V2.3 migration: %s", e)


def _auto_seed_industries(conn):
    """Sync industries and keywords from config/industries.json into DB.

    Uses upsert logic — new industries are added, existing ones are updated.
    This runs on every startup so config file changes are always reflected.
    """
    try:
        config = get_industries()
        config_industries = config.get("industries", [])
        if not config_industries:
            print("[STARTUP] No industries found in config/industries.json!")
            return

        existing_ids = {
            row[0] for row in conn.execute("SELECT id FROM industries").fetchall()
        }
        print(f"[STARTUP] Config has {len(config_industries)} industries, DB has {len(existing_ids)}")
        added = 0
        for ind in config_industries:
            try:
                upsert_industry(conn, ind)
                if ind["id"] not in existing_ids:
                    added += 1
                    print(f"[STARTUP]   + Added industry: {ind['id']} ({ind['name']})")
            except Exception as e:
                print(f"[STARTUP]   ! Failed to upsert industry {ind['id']}: {e}")
            # Seed keywords for this industry
            for kw in ind.get("keywords", []):
                try:
                    conn.execute(
                        "INSERT OR IGNORE INTO keywords (keyword, industry_id, source) "
                        "VALUES (?, ?, 'seeded')",
                        (kw.lower(), ind["id"]),
                    )
                except Exception:
                    pass
        conn.commit()
        total = conn.execute("SELECT COUNT(*) FROM industries").fetchone()[0]
        kw_count = conn.execute("SELECT COUNT(*) FROM keywords").fetchone()[0]
        print(f"[STARTUP] Sync complete: {added} new industries, {total} total, {kw_count} keywords")
    except Exception as e:
        print(f"[STARTUP] Industry sync ERROR: {e}")
        import traceback
        traceback.print_exc()
        logger.warning("Could not sync industries: %s", e)


# ── Pydantic Models ──────────────────────────────────────────────────

class RecipientCreate(BaseModel):
    name: str
    email: str
    role: str = ""
    groups: list[str] = ["executive-team"]
    bu_filter: list[str] = ["all"]
    signal_type_filter: list[str] = ["all"]
    notes: str = ""


class RecipientUpdate(BaseModel):
    name: str | None = None
    email: str | None = None
    role: str | None = None
    groups: list[str] | None = None
    bu_filter: list[str] | None = None
    signal_type_filter: list[str] | None = None
    status: str | None = None
    notes: str | None = None


class SourceCreate(BaseModel):
    name: str
    url: str
    type: str = "rss"
    tier: int = 2
    relevant_bus: list[str] = []
    keywords: list[str] = []
    is_competitor: bool = False


class SourceUpdate(BaseModel):
    name: str | None = None
    url: str | None = None
    active: bool | None = None
    tier: int | None = None
    type: str | None = None
    relevant_bus: list[str] | None = None
    keywords: list[str] | None = None
    is_competitor: bool | None = None


class DeliverySettingsUpdate(BaseModel):
    send_day: str | None = None
    send_time_et: str | None = None
    timezone: str | None = None
    max_recipients_per_batch: int | None = None
    batch_delay_seconds: int | None = None


class IndustryCreate(BaseModel):
    id: str
    name: str
    category: str = ""
    description: str = ""
    relevant_bus: list[str] = []
    keywords: list[str] = []
    priority: int = 2
    active: bool = True


class IndustryUpdate(BaseModel):
    name: str | None = None
    category: str | None = None
    description: str | None = None
    relevant_bus: list[str] | None = None
    priority: int | None = None
    active: bool | None = None


class KeywordCreate(BaseModel):
    keyword: str
    industry_id: str | None = None
    bu_id: str | None = None
    source: str = "manual"


class KeywordBulkImport(BaseModel):
    keywords: list[str]
    industry_id: str | None = None
    source: str = "imported"


# ── Recipients ───────────────────────────────────────────────────────

@app.get("/api/recipients")
def list_recipients():
    """Get all recipients and delivery settings."""
    return get_recipients()


@app.post("/api/recipients")
def add_recipient(recipient: RecipientCreate):
    """Add a new recipient."""
    config = get_recipients()
    recipients = config.get("recipients", [])

    # Check for duplicate email
    for r in recipients:
        if r["email"] == recipient.email:
            raise HTTPException(400, f"Recipient with email {recipient.email} already exists")

    new_id = f"recipient-{len(recipients) + 1:03d}"
    new_recipient = {
        "id": new_id,
        "name": recipient.name,
        "email": recipient.email,
        "role": recipient.role,
        "groups": recipient.groups,
        "bu_filter": recipient.bu_filter,
        "signal_type_filter": recipient.signal_type_filter,
        "status": "active",
        "created_at": datetime.now().isoformat() + "Z",
        "notes": recipient.notes,
    }
    recipients.append(new_recipient)
    config["recipients"] = recipients
    save_recipients(config)

    return new_recipient


@app.put("/api/recipients/{recipient_id}")
def update_recipient(recipient_id: str, updates: RecipientUpdate):
    """Update an existing recipient."""
    config = get_recipients()
    recipients = config.get("recipients", [])

    for r in recipients:
        if r["id"] == recipient_id:
            update_data = updates.model_dump(exclude_none=True)
            r.update(update_data)
            save_recipients(config)
            return r

    raise HTTPException(404, f"Recipient {recipient_id} not found")


@app.delete("/api/recipients/{recipient_id}")
def delete_recipient(recipient_id: str):
    """Remove a recipient."""
    config = get_recipients()
    recipients = config.get("recipients", [])
    original_len = len(recipients)

    config["recipients"] = [r for r in recipients if r["id"] != recipient_id]

    if len(config["recipients"]) == original_len:
        raise HTTPException(404, f"Recipient {recipient_id} not found")

    save_recipients(config)
    return {"deleted": recipient_id}


@app.get("/api/delivery-settings")
def get_delivery_settings():
    """Get the delivery schedule settings."""
    config = get_recipients()
    return config.get("delivery_settings", {})


@app.put("/api/delivery-settings")
def update_delivery_settings(settings: DeliverySettingsUpdate):
    """Update the delivery schedule settings."""
    config = get_recipients()
    current = config.get("delivery_settings", {})
    update_data = settings.model_dump(exclude_none=True)
    current.update(update_data)
    config["delivery_settings"] = current
    save_recipients(config)
    return current


# ── Business Units ───────────────────────────────────────────────────

@app.get("/api/business-units")
def list_business_units():
    """Get all business units configuration."""
    return get_business_units()


@app.put("/api/business-units/{bu_id}")
def update_business_unit(bu_id: str, updates: dict):
    """Update a business unit's configuration."""
    config = get_business_units()
    for bu in config.get("business_units", []):
        if bu["id"] == bu_id:
            bu.update(updates)
            save_business_units(config)
            return bu
    raise HTTPException(404, f"Business unit {bu_id} not found")


# ── Industries (V2.1) ────────────────────────────────────────────────

@app.get("/api/industries")
def list_industries():
    """Get all industries with BU associations and keyword counts."""
    conn = get_connection()
    try:
        industries = get_all_industries(conn)
        return {"industries": industries}
    except Exception:
        # Table may not exist yet — return from config file
        config = get_industries()
        return {"industries": config.get("industries", [])}
    finally:
        conn.close()


@app.post("/api/industries")
def add_industry(industry: IndustryCreate):
    """Add a new industry to both the DB and config."""
    conn = get_connection()
    try:
        upsert_industry(conn, industry.model_dump())
        # Seed keywords if provided
        if industry.keywords:
            bulk_import_keywords(conn, industry.keywords, industry.id)
        # Sync to config file
        _sync_industries_to_config(conn)
        return {"id": industry.id, "status": "created"}
    finally:
        conn.close()


@app.put("/api/industries/{industry_id}")
def update_industry(industry_id: str, updates: IndustryUpdate):
    """Update an existing industry."""
    conn = get_connection()
    try:
        # Fetch existing
        existing = conn.execute("SELECT * FROM industries WHERE id = ?", (industry_id,)).fetchone()
        if not existing:
            raise HTTPException(404, f"Industry {industry_id} not found")
        merged = dict(existing)
        merged["id"] = industry_id
        update_data = updates.model_dump(exclude_none=True)
        merged.update(update_data)
        # Handle active as int for DB
        if "active" in merged and isinstance(merged["active"], bool):
            pass  # upsert_industry handles the conversion
        upsert_industry(conn, merged)
        _sync_industries_to_config(conn)
        return get_all_industries(conn)
    finally:
        conn.close()


@app.delete("/api/industries/{industry_id}")
def remove_industry(industry_id: str):
    """Delete an industry."""
    conn = get_connection()
    try:
        if not db_delete_industry(conn, industry_id):
            raise HTTPException(404, f"Industry {industry_id} not found")
        _sync_industries_to_config(conn)
        return {"deleted": industry_id}
    finally:
        conn.close()


def _sync_industries_to_config(conn):
    """Write current DB industries back to config/industries.json."""
    industries = get_all_industries(conn)
    # Also fetch keywords per industry for config sync
    for ind in industries:
        kws = get_all_keywords(conn, industry_id=ind["id"])
        ind["keywords"] = [k["keyword"] for k in kws]
        # Clean DB-only fields for config
        for key in ("keyword_count", "created_at", "updated_at"):
            ind.pop(key, None)
        ind["active"] = bool(ind.get("active", 1))
    config = get_industries()
    config["industries"] = industries
    save_industries(config)


# ── Keywords (V2.1) ─────────────────────────────────────────────────

@app.get("/api/keywords")
def list_keywords(industry_id: str | None = None, bu_id: str | None = None):
    """Get keywords, optionally filtered by industry or BU."""
    conn = get_connection()
    try:
        keywords = get_all_keywords(conn, industry_id=industry_id, bu_id=bu_id, active_only=False)
        return {"keywords": keywords}
    finally:
        conn.close()


@app.post("/api/keywords")
def add_keyword(kw: KeywordCreate):
    """Add a single keyword."""
    conn = get_connection()
    try:
        kw_id = upsert_keyword(conn, kw.model_dump())
        return {"id": kw_id, "keyword": kw.keyword}
    finally:
        conn.close()


@app.post("/api/keywords/bulk")
def import_keywords_bulk(data: KeywordBulkImport):
    """Bulk import keywords for an industry."""
    conn = get_connection()
    try:
        count = bulk_import_keywords(conn, data.keywords, data.industry_id, data.source)
        return {"imported": count}
    finally:
        conn.close()


@app.delete("/api/keywords/{keyword_id}")
def remove_keyword(keyword_id: int):
    """Delete a keyword."""
    conn = get_connection()
    try:
        if not db_delete_keyword(conn, keyword_id):
            raise HTTPException(404, f"Keyword {keyword_id} not found")
        return {"deleted": keyword_id}
    finally:
        conn.close()


# ── Keyword Discovery (V2.1 Phase B) ────────────────────────────────

@app.get("/api/keywords/discover")
def discover_keywords(min_score: float = 6.0, max_new: int = 20):
    """Discover new keyword candidates from scored signals (preview only)."""
    from src.analyzer.keyword_discovery import discover_keywords_from_signals
    conn = get_connection()
    try:
        return discover_keywords_from_signals(conn, min_score=min_score, max_new=max_new)
    finally:
        conn.close()


@app.post("/api/keywords/discover/import")
def import_discovered_keywords(min_score: float = 6.0, max_new: int = 20, auto_activate: bool = False):
    """Discover and import new keywords into the database."""
    from src.analyzer.keyword_discovery import auto_import_discovered
    conn = get_connection()
    try:
        return auto_import_discovered(conn, min_score=min_score, max_new=max_new, auto_activate=auto_activate)
    finally:
        conn.close()


@app.post("/api/keywords/update-hits")
def update_keyword_hits():
    """Update keyword hit counts based on recent signal matches."""
    from src.analyzer.keyword_discovery import update_keyword_hit_counts
    conn = get_connection()
    try:
        updated = update_keyword_hit_counts(conn)
        return {"updated": updated}
    finally:
        conn.close()


# ── Google Trends (V2.1 Phase B) ────────────────────────────────────

@app.get("/api/google-trends")
def google_trends_snapshot(keywords: str | None = None, timeframe: str = "now 7-d"):
    """Get a Google Trends snapshot for specified keywords.

    Args:
        keywords: Comma-separated keywords (max 5). Uses industry defaults if empty.
        timeframe: Pytrends timeframe string (e.g., 'now 7-d', 'today 3-m').
    """
    from src.collector.trends_collector import get_trend_snapshot
    kw_list = [k.strip() for k in keywords.split(",") if k.strip()][:5] if keywords else None
    return get_trend_snapshot(keywords=kw_list, timeframe=timeframe)


# ── Sources ──────────────────────────────────────────────────────────

@app.get("/api/sources")
def list_sources():
    """Get all configured data sources."""
    return get_sources()


@app.post("/api/sources")
def add_source(source: SourceCreate):
    """Add a new data source."""
    config = get_sources()
    sources = config.get("sources", [])

    # Generate a URL-safe ID from the name
    slug = source.name.lower().replace(" ", "-").replace(".", "")[:30]
    source_id = slug
    existing_ids = {s["id"] for s in sources}
    counter = 1
    while source_id in existing_ids:
        source_id = f"{slug}-{counter}"
        counter += 1

    new_source = {
        "id": source_id,
        "name": source.name,
        "url": source.url,
        "type": source.type,
        "tier": source.tier,
        "relevant_bus": source.relevant_bus,
        "keywords": source.keywords,
        "active": True,
        "last_check": None,
        "error_count": 0,
    }
    if source.is_competitor:
        new_source["is_competitor"] = True

    sources.append(new_source)
    config["sources"] = sources
    save_sources(config)

    return new_source


@app.put("/api/sources/{source_id}")
def update_source(source_id: str, updates: SourceUpdate):
    """Update a data source."""
    config = get_sources()
    for source in config.get("sources", []):
        if source["id"] == source_id:
            update_data = updates.model_dump(exclude_none=True)
            source.update(update_data)
            save_sources(config)
            return source
    raise HTTPException(404, f"Source {source_id} not found")


@app.delete("/api/sources/{source_id}")
def delete_source(source_id: str):
    """Remove a data source."""
    config = get_sources()
    sources = config.get("sources", [])
    original_len = len(sources)

    config["sources"] = [s for s in sources if s["id"] != source_id]

    if len(config["sources"]) == original_len:
        raise HTTPException(404, f"Source {source_id} not found")

    save_sources(config)
    return {"deleted": source_id}


# ── Trends ───────────────────────────────────────────────────────────

@app.get("/api/trends")
def list_trends(limit: int = 30, start_date: str | None = None, end_date: str | None = None):
    """Get current trend summary, optionally filtered by date range."""
    from src.trends.tracker import get_trend_summary
    return get_trend_summary(limit=limit, start_date=start_date, end_date=end_date)


@app.get("/api/trends/sentiment")
def trends_sentiment():
    """Get sentiment/momentum breakdown per BU and per Industry.

    Returns aggregated signal counts, average scores, and momentum
    indicators grouped by Business Unit and Industry.
    """
    conn = get_connection()
    try:
        bu_config = get_business_units()
        bu_names = {bu["id"]: bu.get("short_name", bu["name"])
                    for bu in bu_config.get("business_units", [])}

        # BU sentiment: count, avg score, momentum per BU
        bu_rows = conn.execute("""
            SELECT sb.bu_id, COUNT(*) as cnt,
                   AVG(sa.score_composite) as avg_score,
                   SUM(CASE WHEN sa.signal_type = 'competitive-threat' THEN 1 ELSE 0 END) as threats,
                   SUM(CASE WHEN sa.signal_type = 'revenue-opportunity' THEN 1 ELSE 0 END) as opps,
                   SUM(CASE WHEN sa.signal_type IN ('technology-trend', 'market-shift') THEN 1 ELSE 0 END) as shifts
            FROM signal_bus sb
            JOIN signal_analysis sa ON sb.signal_id = sa.signal_id
            JOIN signals s ON sb.signal_id = s.id
            WHERE s.status IN ('scored', 'published')
              AND COALESCE(s.dismissed, 0) = 0
            GROUP BY sb.bu_id
            ORDER BY avg_score DESC
        """).fetchall()

        bu_sentiment = []
        for r in bu_rows:
            total = r[1]
            threats = r[3]
            opps = r[4]
            sentiment_score = (opps - threats) / max(total, 1)
            if sentiment_score > 0.2:
                sentiment = "positive"
            elif sentiment_score < -0.2:
                sentiment = "negative"
            else:
                sentiment = "neutral"

            bu_sentiment.append({
                "id": r[0],
                "name": bu_names.get(r[0], r[0]),
                "signal_count": total,
                "avg_score": round(r[2] or 0, 1),
                "threats": threats,
                "opportunities": opps,
                "shifts": r[5],
                "sentiment": sentiment,
                "sentiment_score": round(sentiment_score, 2),
            })

        # Industry sentiment
        ind_rows = conn.execute("""
            SELECT si.industry_id, i.name, COUNT(*) as cnt,
                   AVG(sa.score_composite) as avg_score,
                   SUM(CASE WHEN sa.signal_type = 'competitive-threat' THEN 1 ELSE 0 END) as threats,
                   SUM(CASE WHEN sa.signal_type = 'revenue-opportunity' THEN 1 ELSE 0 END) as opps
            FROM signal_industries si
            JOIN industries i ON si.industry_id = i.id
            JOIN signal_analysis sa ON si.signal_id = sa.signal_id
            JOIN signals s ON si.signal_id = s.id
            WHERE s.status IN ('scored', 'published')
              AND COALESCE(s.dismissed, 0) = 0
            GROUP BY si.industry_id
            ORDER BY avg_score DESC
        """).fetchall()

        ind_sentiment = []
        for r in ind_rows:
            total = r[2]
            threats = r[4]
            opps = r[5]
            sentiment_score = (opps - threats) / max(total, 1)
            if sentiment_score > 0.2:
                sentiment = "positive"
            elif sentiment_score < -0.2:
                sentiment = "negative"
            else:
                sentiment = "neutral"

            ind_sentiment.append({
                "id": r[0],
                "name": r[1],
                "signal_count": total,
                "avg_score": round(r[3] or 0, 1),
                "threats": threats,
                "opportunities": opps,
                "sentiment": sentiment,
                "sentiment_score": round(sentiment_score, 2),
            })

        return {
            "bu_sentiment": bu_sentiment,
            "industry_sentiment": ind_sentiment,
        }
    finally:
        conn.close()


@app.get("/api/trends/{trend_key:path}/history")
def trend_history(trend_key: str, weeks: int = 12):
    """Get week-by-week history for a specific trend."""
    from src.trends.tracker import get_trend_history
    return {"trend_key": trend_key, "history": get_trend_history(trend_key, weeks)}


@app.get("/api/trends/alerts")
def list_trend_alerts(bu_code: str | None = None, industry: str | None = None,
                      limit: int = 5):
    """Get AI-generated 'What's Moving' trend alert cards."""
    from src.trends.tracker import get_trend_alerts
    return {"alerts": get_trend_alerts(bu_code=bu_code, industry=industry, limit=limit)}


@app.post("/api/trends/alerts/generate")
def generate_alerts(bu_code: str | None = None, start_date: str | None = None,
                    end_date: str | None = None):
    """Generate new trend alerts from current signal data."""
    from src.trends.tracker import generate_trend_alerts
    alerts = generate_trend_alerts(bu_code=bu_code, start_date=start_date, end_date=end_date)
    return {"generated": len(alerts), "alerts": alerts}


@app.get("/api/trends/industry-momentum")
def industry_momentum(bu_code: str | None = None, start_date: str | None = None,
                      end_date: str | None = None):
    """Get industry momentum cards with sparkline data, sentiment, and scores."""
    from src.trends.tracker import get_industry_momentum
    return {"industries": get_industry_momentum(bu_code=bu_code, start_date=start_date,
                                                 end_date=end_date)}


@app.get("/api/trends/signal-volume")
def signal_volume_over_time(bu_code: str | None = None, weeks: int = 12):
    """Get weekly signal volume data for multi-line chart."""
    from src.trends.tracker import get_signal_volume_over_time
    return get_signal_volume_over_time(bu_code=bu_code, weeks=weeks)


@app.get("/api/trends/competitor-trends")
def competitor_trends(bu_code: str | None = None, start_date: str | None = None,
                      end_date: str | None = None):
    """Get competitor signal trends with period-over-period comparison."""
    from src.trends.tracker import get_competitor_trends
    return {"competitors": get_competitor_trends(bu_code=bu_code, start_date=start_date,
                                                  end_date=end_date)}


@app.get("/api/export/trends/excel")
def export_trends_excel(bu_code: str | None = None, start_date: str | None = None,
                        end_date: str | None = None):
    """Export Trends data as Excel workbook with specialized tabs."""
    try:
        import openpyxl  # noqa: F401
    except ImportError:
        raise HTTPException(422, detail="Excel export requires openpyxl. pip install openpyxl")
    from src.export.trends_export import export_trends_excel as _export
    from fastapi.responses import StreamingResponse
    buffer = _export(bu_code=bu_code, start_date=start_date, end_date=end_date)
    filename = f"vpg-trends-{datetime.now().strftime('%Y%m%d')}.xlsx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@app.get("/api/export/trends/pptx")
def export_trends_pptx(bu_code: str | None = None, start_date: str | None = None,
                       end_date: str | None = None):
    """Export Trends data as PowerPoint presentation."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        raise HTTPException(422, detail="PowerPoint export requires python-pptx. pip install python-pptx")
    from src.export.trends_export import export_trends_pptx as _export
    from fastapi.responses import StreamingResponse
    buffer = _export(bu_code=bu_code, start_date=start_date, end_date=end_date)
    filename = f"vpg-trends-{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


# ── Scoring ──────────────────────────────────────────────────────────

@app.get("/api/scoring")
def get_scoring():
    """Get scoring configuration."""
    return get_scoring_weights()


@app.put("/api/scoring")
def update_scoring(weights: dict):
    """Update scoring weights and thresholds."""
    save_scoring_weights(weights)
    return get_scoring_weights()


# ── Export (V2.1 Phase C) ────────────────────────────────────────────

@app.get("/api/export/check")
def export_check():
    """Check which export formats are available (have required packages installed)."""
    excel_ok = False
    pptx_ok = False
    try:
        import openpyxl  # noqa: F401
        excel_ok = True
    except ImportError:
        pass
    try:
        import pptx  # noqa: F401
        pptx_ok = True
    except ImportError:
        pass
    return {
        "excel": {"available": excel_ok, "install": "pip install openpyxl"},
        "pptx": {"available": pptx_ok, "install": "pip install python-pptx"},
    }


@app.get("/api/export/excel")
def export_excel(start_date: str | None = None, end_date: str | None = None,
                 bu_id: str | None = None, signal_type: str | None = None,
                 industry_id: str | None = None, min_score: float = 0):
    """Export signals, trends, and keywords as an Excel workbook with applied filters."""
    try:
        import openpyxl  # noqa: F401
    except ImportError:
        raise HTTPException(
            422,
            detail="Excel export requires the 'openpyxl' package. "
                   "Install it with: pip install openpyxl"
        )
    from src.export.excel_export import export_signals_excel
    from fastapi.responses import StreamingResponse
    buffer = export_signals_excel(
        start_date=start_date, end_date=end_date,
        bu_id=bu_id, signal_type=signal_type,
        industry_id=industry_id, min_score=min_score,
    )
    filename = f"vpg-intel-{datetime.now().strftime('%Y%m%d')}.xlsx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@app.get("/api/export/pptx")
def export_pptx(start_date: str | None = None, end_date: str | None = None,
                bu_id: str | None = None, signal_type: str | None = None,
                industry_id: str | None = None, min_score: float = 0,
                max_signals: int = 10):
    """Export intelligence as a PowerPoint presentation with applied filters."""
    try:
        import pptx  # noqa: F401
    except ImportError:
        raise HTTPException(
            422,
            detail="PowerPoint export requires the 'python-pptx' package. "
                   "Install it with: pip install python-pptx"
        )
    from src.export.pptx_export import export_signals_pptx
    from fastapi.responses import StreamingResponse
    buffer = export_signals_pptx(
        start_date=start_date, end_date=end_date,
        bu_id=bu_id, signal_type=signal_type,
        industry_id=industry_id, min_score=min_score,
        max_signals=max_signals,
    )
    filename = f"vpg-intel-{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


# ── Intelligence Feed (V2.1 Phase C) ────────────────────────────────

@app.get("/api/feed")
def intelligence_feed(start_date: str | None = None, end_date: str | None = None,
                      bu_id: str | None = None, signal_type: str | None = None,
                      industry_id: str | None = None, min_score: float = 0,
                      limit: int = 50):
    """Get a filterable feed of scored signals with full analysis data.

    Powers the Intelligence Feed UI component.
    """
    conn = get_connection()
    try:
        query = """
            SELECT s.id, s.title, s.url, s.source_name, s.source_tier,
                   s.published_at, s.collected_at, s.image_url,
                   sa.signal_type, sa.headline, sa.what_summary, sa.why_it_matters,
                   sa.quick_win, sa.suggested_owner, sa.estimated_impact,
                   sa.score_composite, sa.score_revenue_impact, sa.score_time_sensitivity,
                   sa.score_strategic_alignment, sa.score_competitive_pressure,
                   sa.validation_level,
                   COALESCE(s.dismissed, 0), COALESCE(s.handled, 0)
            FROM signals s
            JOIN signal_analysis sa ON s.id = sa.signal_id
            WHERE s.status IN ('scored', 'published')
              AND COALESCE(s.dismissed, 0) = 0
        """
        params: list = []

        if min_score > 0:
            query += " AND sa.score_composite >= ?"
            params.append(min_score)
        if start_date:
            query += " AND COALESCE(s.published_at, s.collected_at) >= ?"
            params.append(start_date)
        if end_date:
            query += " AND COALESCE(s.published_at, s.collected_at) <= ?"
            params.append(end_date + " 23:59:59")
        if signal_type:
            query += " AND sa.signal_type = ?"
            params.append(signal_type)
        if bu_id:
            query += " AND s.id IN (SELECT signal_id FROM signal_bus WHERE bu_id = ?)"
            params.append(bu_id)
        if industry_id:
            query += " AND s.id IN (SELECT signal_id FROM signal_industries WHERE industry_id = ?)"
            params.append(industry_id)

        query += " ORDER BY sa.score_composite DESC LIMIT ?"
        params.append(limit)

        rows = conn.execute(query, params).fetchall()

        # Also fetch source links for each signal (validations)
        signals = []
        for row in rows:
            sig_id = row[0]
            bus = conn.execute(
                "SELECT bu_id FROM signal_bus WHERE signal_id = ?", (sig_id,)
            ).fetchall()
            inds = conn.execute(
                "SELECT industry_id FROM signal_industries WHERE signal_id = ?", (sig_id,)
            ).fetchall()
            sources = conn.execute(
                "SELECT corroborating_url, corroborating_source, corroborating_title "
                "FROM signal_validations WHERE signal_id = ?", (sig_id,)
            ).fetchall()

            signals.append({
                "id": sig_id, "title": row[1], "url": row[2],
                "source_name": row[3], "source_tier": row[4],
                "published_at": row[5], "collected_at": row[6], "image_url": row[7],
                "signal_type": row[8], "headline": row[9],
                "what_summary": row[10], "why_it_matters": row[11],
                "quick_win": row[12], "suggested_owner": row[13],
                "estimated_impact": row[14],
                "score_composite": row[15], "score_revenue_impact": row[16],
                "score_time_sensitivity": row[17], "score_strategic_alignment": row[18],
                "score_competitive_pressure": row[19],
                "validation_level": row[20],
                "bus": [b[0] for b in bus],
                "industries": [i[0] for i in inds],
                "source_links": [{"url": s[0], "source": s[1], "title": s[2]} for s in sources],
                "dismissed": row[21],
                "handled": row[22],
            })

        return {"signals": signals, "total": len(signals)}
    finally:
        conn.close()


# ── BU Executive Dashboard Data (V2.1 Phase C) ──────────────────────

@app.get("/api/executive/bu-summary")
def bu_executive_summary():
    """Get per-BU signal summary for the executive dashboard."""
    conn = get_connection()
    try:
        bu_config = get_business_units()
        bu_names = {bu["id"]: bu for bu in bu_config.get("business_units", [])}

        # Get signal counts and scores per BU
        rows = conn.execute("""
            SELECT sb.bu_id, COUNT(*) as cnt,
                   AVG(sa.score_composite) as avg_score,
                   MAX(sa.score_composite) as max_score
            FROM signal_bus sb
            JOIN signal_analysis sa ON sb.signal_id = sa.signal_id
            JOIN signals s ON sb.signal_id = s.id
            WHERE s.status IN ('scored', 'published')
            GROUP BY sb.bu_id
            ORDER BY avg_score DESC
        """).fetchall()

        summaries = []
        for row in rows:
            bu_id = row[0]
            bu_info = bu_names.get(bu_id, {})

            # Top signal types for this BU
            type_rows = conn.execute("""
                SELECT sa.signal_type, COUNT(*) as cnt
                FROM signal_bus sb
                JOIN signal_analysis sa ON sb.signal_id = sa.signal_id
                JOIN signals s ON sb.signal_id = s.id
                WHERE sb.bu_id = ? AND s.status IN ('scored', 'published')
                GROUP BY sa.signal_type ORDER BY cnt DESC LIMIT 3
            """, (bu_id,)).fetchall()

            # Top signal for this BU
            top_signal = conn.execute("""
                SELECT sa.headline, sa.score_composite
                FROM signal_bus sb
                JOIN signal_analysis sa ON sb.signal_id = sa.signal_id
                JOIN signals s ON sb.signal_id = s.id
                WHERE sb.bu_id = ? AND s.status IN ('scored', 'published')
                ORDER BY sa.score_composite DESC LIMIT 1
            """, (bu_id,)).fetchone()

            summaries.append({
                "bu_id": bu_id,
                "bu_name": bu_info.get("name", bu_id),
                "bu_short": bu_info.get("short_name", bu_id),
                "color": bu_info.get("color", "#2E75B6"),
                "signal_count": row[1],
                "avg_score": round(row[2] or 0, 1),
                "max_score": round(row[3] or 0, 1),
                "top_types": [{"type": t[0], "count": t[1]} for t in type_rows],
                "top_signal": {
                    "headline": top_signal[0], "score": round(top_signal[1] or 0, 1)
                } if top_signal else None,
            })

        # Overall stats
        total_signals = conn.execute(
            "SELECT COUNT(*) FROM signals WHERE status IN ('scored', 'published')"
        ).fetchone()[0]
        avg_all = conn.execute(
            "SELECT AVG(score_composite) FROM signal_analysis"
        ).fetchone()[0]

        return {
            "bu_summaries": summaries,
            "total_signals": total_signals,
            "overall_avg_score": round(avg_all or 0, 1),
            "bus_with_signals": len(summaries),
            "generated_at": datetime.now().isoformat(),
        }
    finally:
        conn.close()


# ── Executive Dashboard V2.4 ─────────────────────────────────────────

@app.get("/api/executive/dashboard")
def executive_dashboard(bu_id: str | None = None, industry_id: str | None = None,
                        start_date: str | None = None, end_date: str | None = None):
    """V2.4 Executive Dashboard — All-BU heatmap or Single-BU deep view.

    Returns all data needed by the redesigned Executive Dashboard in one call.
    """
    from src.analyzer.recommendations import generate_recommendations

    conn = get_connection()
    try:
        bu_config = get_business_units()
        all_bus = {bu["id"]: bu for bu in bu_config.get("business_units", []) if bu.get("active", True)}
        ind_config = get_industries()
        all_industries = {ind["id"]: ind for ind in ind_config.get("industries", [])}

        # Date filter helpers
        date_where = ""
        date_params: list = []
        if start_date:
            date_where += " AND COALESCE(s.published_at, s.collected_at) >= ?"
            date_params.append(start_date)
        if end_date:
            date_where += " AND COALESCE(s.published_at, s.collected_at) <= ?"
            date_params.append(end_date + " 23:59:59")

        mode = "single-bu" if bu_id else "all-bu"
        result: dict = {"mode": mode, "generated_at": datetime.now().isoformat()}

        # ── BU Heatmap (All-BU mode) ────────────────────────────
        if mode == "all-bu":
            heatmap = []
            for bid, binfo in all_bus.items():
                # Signal count this period
                q = """SELECT COUNT(*) FROM signals s
                       JOIN signal_bus sb ON s.id = sb.signal_id
                       JOIN signal_analysis sa ON s.id = sa.signal_id
                       WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                         AND COALESCE(s.dismissed,0) = 0""" + date_where
                cnt = conn.execute(q, [bid] + date_params).fetchone()[0]

                # Prior period count (same duration shifted back)
                prior_cnt = 0
                if start_date and end_date:
                    from datetime import timedelta
                    d_start = datetime.strptime(start_date, "%Y-%m-%d")
                    d_end = datetime.strptime(end_date, "%Y-%m-%d")
                    delta = (d_end - d_start).days
                    p_end = (d_start - timedelta(days=1)).strftime("%Y-%m-%d")
                    p_start = (d_start - timedelta(days=delta + 1)).strftime("%Y-%m-%d")
                    q2 = """SELECT COUNT(*) FROM signals s
                            JOIN signal_bus sb ON s.id = sb.signal_id
                            JOIN signal_analysis sa ON s.id = sa.signal_id
                            WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                              AND COALESCE(s.dismissed,0) = 0
                              AND COALESCE(s.published_at, s.collected_at) >= ?
                              AND COALESCE(s.published_at, s.collected_at) <= ?"""
                    prior_cnt = conn.execute(q2, [bid, p_start, p_end + " 23:59:59"]).fetchone()[0]
                else:
                    # Default: compare last 14 days vs prior 14 days
                    q_recent = """SELECT COUNT(*) FROM signals s
                                  JOIN signal_bus sb ON s.id = sb.signal_id
                                  WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                                    AND COALESCE(s.dismissed,0) = 0
                                    AND s.collected_at >= date('now', '-14 days')"""
                    q_prior = """SELECT COUNT(*) FROM signals s
                                 JOIN signal_bus sb ON s.id = sb.signal_id
                                 WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                                   AND COALESCE(s.dismissed,0) = 0
                                   AND s.collected_at >= date('now', '-28 days')
                                   AND s.collected_at < date('now', '-14 days')"""
                    cnt_r = conn.execute(q_recent, [bid]).fetchone()[0]
                    prior_cnt = conn.execute(q_prior, [bid]).fetchone()[0]
                    if not start_date and not end_date:
                        # Use total for display, recent for trend
                        pass

                # Trend arrow
                if cnt > prior_cnt:
                    trend = "up"
                elif cnt < prior_cnt:
                    trend = "down"
                else:
                    trend = "stable"

                trend_pct = 0
                if prior_cnt > 0:
                    trend_pct = round(((cnt - prior_cnt) / prior_cnt) * 100)

                # Top signal type
                top_type_row = conn.execute("""
                    SELECT sa.signal_type, COUNT(*) as c FROM signal_bus sb
                    JOIN signal_analysis sa ON sb.signal_id = sa.signal_id
                    JOIN signals s ON sb.signal_id = s.id
                    WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                      AND COALESCE(s.dismissed,0) = 0
                    GROUP BY sa.signal_type ORDER BY c DESC LIMIT 1
                """, (bid,)).fetchone()

                heatmap.append({
                    "bu_id": bid,
                    "bu_name": binfo.get("name", bid),
                    "bu_short": binfo.get("short_name", bid),
                    "color": binfo.get("color", "#2E75B6"),
                    "signal_count": cnt,
                    "trend": trend,
                    "trend_pct": trend_pct,
                    "top_signal_type": top_type_row[0] if top_type_row else None,
                    "has_critical": False,  # Set below after recs
                })

            result["bu_heatmap"] = heatmap

        # ── Top 5 Actions (score >= 7.0) ────────────────────────
        actions_q = """
            SELECT s.id, sa.headline, sa.signal_type, sa.score_composite,
                   sa.what_summary, sa.quick_win, sa.suggested_owner,
                   sa.estimated_impact, sa.validation_level
            FROM signals s
            JOIN signal_analysis sa ON s.id = sa.signal_id
            WHERE s.status IN ('scored','published')
              AND COALESCE(s.dismissed,0) = 0
              AND sa.score_composite >= 7.0
        """ + date_where
        actions_params = list(date_params)
        if bu_id:
            actions_q += " AND s.id IN (SELECT signal_id FROM signal_bus WHERE bu_id = ?)"
            actions_params.append(bu_id)
        if industry_id:
            actions_q += " AND s.id IN (SELECT signal_id FROM signal_industries WHERE industry_id = ?)"
            actions_params.append(industry_id)
        actions_q += " ORDER BY sa.score_composite DESC LIMIT 5"

        action_rows = conn.execute(actions_q, actions_params).fetchall()
        top_actions = []
        for r in action_rows:
            sig_id = r[0]
            bus = [b[0] for b in conn.execute(
                "SELECT bu_id FROM signal_bus WHERE signal_id = ?", (sig_id,)
            ).fetchall()]
            inds = [i[0] for i in conn.execute(
                "SELECT industry_id FROM signal_industries WHERE signal_id = ?", (sig_id,)
            ).fetchall()]
            top_actions.append({
                "id": sig_id, "headline": r[1], "signal_type": r[2],
                "score": round(r[3] or 0, 1), "what_summary": r[4],
                "quick_win": r[5], "owner": r[6], "impact": r[7],
                "validation": r[8], "bus": bus, "industries": inds,
            })
        result["top_actions"] = top_actions

        # ── Alerts (critical/high recs + persistent signals) ────
        try:
            recs_data = generate_recommendations(conn, max_recommendations=20)
            all_recs = recs_data.get("recommendations", [])
        except Exception:
            all_recs = []

        alerts = []
        for rec in all_recs:
            if rec.get("priority", 3) <= 2:  # critical or high
                rec_bu = rec.get("bu_id") or rec.get("bu_ids", "")
                if bu_id and rec_bu and bu_id not in str(rec_bu):
                    continue
                alerts.append({
                    "type": rec.get("type", "recommendation"),
                    "priority": rec.get("priority", 2),
                    "title": rec.get("title", ""),
                    "description": rec.get("description", ""),
                    "action": rec.get("action", ""),
                    "score": rec.get("score", 0),
                    "bu_id": rec_bu,
                })

        result["alerts"] = alerts[:8]

        # Mark BUs with critical alerts in heatmap
        if mode == "all-bu":
            critical_bus = set()
            for a in alerts:
                if a["priority"] == 1 and a.get("bu_id"):
                    critical_bus.add(a["bu_id"])
            for card in result["bu_heatmap"]:
                if card["bu_id"] in critical_bus:
                    card["has_critical"] = True

        # ── Competitor Pulse (enhanced) ──────────────────────────
        comp_results = []
        for comp in MONITORED_COMPETITORS:
            cq = """SELECT COUNT(*), AVG(sa.score_composite)
                    FROM signals s
                    JOIN signal_analysis sa ON s.id = sa.signal_id
                    WHERE s.status IN ('scored','published')
                      AND COALESCE(s.dismissed,0) = 0
                      AND (LOWER(s.title) LIKE ? OR LOWER(s.raw_content) LIKE ?
                           OR LOWER(sa.headline) LIKE ? OR LOWER(sa.what_summary) LIKE ?)
                    """ + date_where
            like = f"%{comp.lower()}%"
            cp = [like, like, like, like] + date_params
            if bu_id:
                cq += " AND s.id IN (SELECT signal_id FROM signal_bus WHERE bu_id = ?)"
                cp.append(bu_id)
            row = conn.execute(cq, cp).fetchone()
            count = row[0] or 0
            if count == 0:
                continue

            # Check if this competitor is new this period
            is_new = False
            if start_date:
                prior_check = conn.execute(
                    "SELECT COUNT(*) FROM signals s JOIN signal_analysis sa ON s.id = sa.signal_id "
                    "WHERE s.status IN ('scored','published') AND COALESCE(s.dismissed,0) = 0 "
                    "AND (LOWER(s.title) LIKE ? OR LOWER(sa.headline) LIKE ?) "
                    "AND COALESCE(s.published_at, s.collected_at) < ?",
                    (like, like, start_date)
                ).fetchone()[0]
                is_new = prior_check == 0

            # Trend
            recent = conn.execute(
                "SELECT COUNT(*) FROM signals s JOIN signal_analysis sa ON s.id = sa.signal_id "
                "WHERE s.status IN ('scored','published') AND COALESCE(s.dismissed,0)=0 "
                "AND (LOWER(s.title) LIKE ? OR LOWER(sa.headline) LIKE ?) "
                "AND s.collected_at >= date('now', '-14 days')",
                (like, like)
            ).fetchone()[0]
            older = conn.execute(
                "SELECT COUNT(*) FROM signals s JOIN signal_analysis sa ON s.id = sa.signal_id "
                "WHERE s.status IN ('scored','published') AND COALESCE(s.dismissed,0)=0 "
                "AND (LOWER(s.title) LIKE ? OR LOWER(sa.headline) LIKE ?) "
                "AND s.collected_at >= date('now', '-28 days') AND s.collected_at < date('now', '-14 days')",
                (like, like)
            ).fetchone()[0]

            trend = "up" if recent > older else ("down" if recent < older else "stable")

            comp_results.append({
                "competitor": comp, "signal_count": count,
                "avg_score": round(row[1] or 0, 1), "trend": trend,
                "is_new": is_new, "recent_signals": recent,
            })

        comp_results.sort(key=lambda x: x["signal_count"], reverse=True)
        result["competitor_pulse"] = comp_results[:10]

        # ── Single-BU Deep View extras ──────────────────────────
        if mode == "single-bu" and bu_id:
            binfo = all_bus.get(bu_id, {})

            # BU header stats
            total_q = """SELECT COUNT(*), AVG(sa.score_composite)
                         FROM signals s JOIN signal_bus sb ON s.id = sb.signal_id
                         JOIN signal_analysis sa ON s.id = sa.signal_id
                         WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                           AND COALESCE(s.dismissed,0) = 0""" + date_where
            hdr = conn.execute(total_q, [bu_id] + date_params).fetchone()

            # Prior period for trend
            if start_date and end_date:
                from datetime import timedelta
                d_start = datetime.strptime(start_date, "%Y-%m-%d")
                d_end = datetime.strptime(end_date, "%Y-%m-%d")
                delta = (d_end - d_start).days
                p_end = (d_start - timedelta(days=1)).strftime("%Y-%m-%d")
                p_start = (d_start - timedelta(days=delta + 1)).strftime("%Y-%m-%d")
                prior_total = conn.execute(
                    """SELECT COUNT(*) FROM signals s JOIN signal_bus sb ON s.id = sb.signal_id
                       WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                         AND COALESCE(s.dismissed,0) = 0
                         AND COALESCE(s.published_at, s.collected_at) >= ?
                         AND COALESCE(s.published_at, s.collected_at) <= ?""",
                    [bu_id, p_start, p_end + " 23:59:59"]
                ).fetchone()[0]
            else:
                prior_total = conn.execute(
                    """SELECT COUNT(*) FROM signals s JOIN signal_bus sb ON s.id = sb.signal_id
                       WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                         AND COALESCE(s.dismissed,0) = 0
                         AND s.collected_at >= date('now', '-28 days')
                         AND s.collected_at < date('now', '-14 days')""",
                    [bu_id]
                ).fetchone()[0]

            current_total = hdr[0] or 0
            trend_pct = round(((current_total - prior_total) / prior_total) * 100) if prior_total > 0 else 0

            result["bu_header"] = {
                "bu_id": bu_id,
                "bu_name": binfo.get("name", bu_id),
                "bu_short": binfo.get("short_name", bu_id),
                "color": binfo.get("color", "#2E75B6"),
                "signal_count": current_total,
                "avg_score": round(hdr[1] or 0, 1),
                "trend_pct": trend_pct,
            }

            # Top industries for this BU
            ind_rows = conn.execute("""
                SELECT si.industry_id, COUNT(*) as cnt, AVG(sa.score_composite) as avg_s
                FROM signal_industries si
                JOIN signal_bus sb ON si.signal_id = sb.signal_id
                JOIN signals s ON si.signal_id = s.id
                JOIN signal_analysis sa ON si.signal_id = sa.signal_id
                WHERE sb.bu_id = ? AND s.status IN ('scored','published')
                  AND COALESCE(s.dismissed,0) = 0
                GROUP BY si.industry_id ORDER BY cnt DESC
            """, (bu_id,)).fetchall()

            industry_breakdown = []
            for ir in ind_rows:
                ind_id = ir[0]
                ind_info = all_industries.get(ind_id, {})

                # Top signal for this industry+BU
                top_sig = conn.execute("""
                    SELECT sa.headline, sa.score_composite FROM signals s
                    JOIN signal_analysis sa ON s.id = sa.signal_id
                    JOIN signal_bus sb ON s.id = sb.signal_id
                    JOIN signal_industries si ON s.id = si.signal_id
                    WHERE sb.bu_id = ? AND si.industry_id = ?
                      AND s.status IN ('scored','published')
                      AND COALESCE(s.dismissed,0) = 0
                    ORDER BY sa.score_composite DESC LIMIT 1
                """, (bu_id, ind_id)).fetchone()

                # Trend for this industry
                ind_recent = conn.execute("""
                    SELECT COUNT(*) FROM signals s
                    JOIN signal_bus sb ON s.id = sb.signal_id
                    JOIN signal_industries si ON s.id = si.signal_id
                    WHERE sb.bu_id = ? AND si.industry_id = ?
                      AND s.status IN ('scored','published')
                      AND COALESCE(s.dismissed,0) = 0
                      AND s.collected_at >= date('now', '-14 days')
                """, (bu_id, ind_id)).fetchone()[0]
                ind_older = conn.execute("""
                    SELECT COUNT(*) FROM signals s
                    JOIN signal_bus sb ON s.id = sb.signal_id
                    JOIN signal_industries si ON s.id = si.signal_id
                    WHERE sb.bu_id = ? AND si.industry_id = ?
                      AND s.status IN ('scored','published')
                      AND COALESCE(s.dismissed,0) = 0
                      AND s.collected_at >= date('now', '-28 days')
                      AND s.collected_at < date('now', '-14 days')
                """, (bu_id, ind_id)).fetchone()[0]

                ind_trend_pct = round(((ind_recent - ind_older) / ind_older) * 100) if ind_older > 0 else 0

                industry_breakdown.append({
                    "industry_id": ind_id,
                    "industry_name": ind_info.get("name", ind_id),
                    "signal_count": ir[1],
                    "avg_score": round(ir[2] or 0, 1),
                    "trend_pct": ind_trend_pct,
                    "trending": "up" if ind_recent > ind_older else ("down" if ind_recent < ind_older else "stable"),
                    "top_signal": top_sig[0] if top_sig else None,
                })

            result["industry_breakdown"] = industry_breakdown

            # Top industries for header
            result["bu_header"]["top_industries"] = [
                {"name": ib["industry_name"], "count": ib["signal_count"]}
                for ib in industry_breakdown[:5]
            ]

            # BU-specific recommendations
            bu_recs = [r for r in all_recs if bu_id in str(r.get("bu_id", "")) or bu_id in str(r.get("bu_ids", ""))]
            result["bu_recommendations"] = [{
                "type": r.get("type", ""),
                "priority": r.get("priority", 3),
                "title": r.get("title", ""),
                "description": r.get("description", ""),
                "action": r.get("action", ""),
                "score": r.get("score", 0),
            } for r in bu_recs[:5]]

        return result

    finally:
        conn.close()


# ── Pipeline Control ─────────────────────────────────────────────────

def _run_pipeline_bg(dry_run: bool = False, pdf_mode: bool = False):
    """Run the pipeline in a background thread."""
    _pipeline_status["running"] = True
    _pipeline_status["paused"] = False
    _pipeline_status["current_stage"] = ""
    _pipeline_status["last_run"] = datetime.now().isoformat()
    try:
        if dry_run:
            from scripts.dry_run import main as dry_run_main
            dry_run_main(pdf_mode=pdf_mode)
            _pipeline_status["last_result"] = {"status": "completed", "mode": "dry-run", "pdf_generated": pdf_mode}
        else:
            from src.pipeline import run_full_pipeline
            result = run_full_pipeline(pdf_mode=pdf_mode)
            _pipeline_status["last_result"] = result
    except Exception as e:
        _pipeline_status["last_result"] = {"status": "failed", "error": str(e)}
    finally:
        _pipeline_status["running"] = False
        _pipeline_status["paused"] = False
        _pipeline_status["current_stage"] = ""


@app.post("/api/pipeline/run")
def run_pipeline(dry_run: bool = False, pdf_mode: bool = False):
    """Trigger a pipeline run (live or dry-run).

    Args:
        dry_run: If true, use seed data instead of live sources.
        pdf_mode: If true, generate PDF and send as attachment.

    Runs asynchronously in the background. Check /api/pipeline/status for results.
    """
    if _pipeline_status["running"]:
        raise HTTPException(409, "Pipeline is already running")

    thread = threading.Thread(
        target=_run_pipeline_bg,
        args=(dry_run, pdf_mode),
        daemon=True,
    )
    thread.start()

    return {
        "message": f"Pipeline started ({'dry-run' if dry_run else 'live'} mode, "
                   f"{'PDF' if pdf_mode else 'HTML'} delivery)",
        "status_url": "/api/pipeline/status",
    }


@app.post("/api/pipeline/pause")
def pause_pipeline():
    """Pause a running pipeline."""
    if not _pipeline_status["running"]:
        raise HTTPException(400, "Pipeline is not running")

    from src.pipeline import pipeline_control
    pipeline_control.pause()
    _pipeline_status["paused"] = True
    return {"message": "Pipeline paused", "paused": True}


@app.post("/api/pipeline/resume")
def resume_pipeline():
    """Resume a paused pipeline."""
    if not _pipeline_status["running"]:
        raise HTTPException(400, "Pipeline is not running")

    from src.pipeline import pipeline_control
    pipeline_control.resume()
    _pipeline_status["paused"] = False
    return {"message": "Pipeline resumed", "paused": False}


@app.post("/api/pipeline/cancel")
def cancel_pipeline():
    """Cancel a running pipeline."""
    if not _pipeline_status["running"]:
        raise HTTPException(400, "Pipeline is not running")

    from src.pipeline import pipeline_control
    pipeline_control.cancel()
    return {"message": "Pipeline cancellation requested"}


@app.get("/api/pipeline/status")
def pipeline_status():
    """Get the current pipeline status."""
    # Sync stage info from the pipeline control
    if _pipeline_status["running"]:
        try:
            from src.pipeline import pipeline_control
            _pipeline_status["current_stage"] = pipeline_control.current_stage
            _pipeline_status["paused"] = pipeline_control.is_paused
        except Exception:
            pass

    return _pipeline_status


# ── Scan / Send Split (V2.4) ────────────────────────────────────────

def _run_scan_bg(sources: list[str] | None, date_from: str | None, date_to: str | None):
    """Run scan-only pipeline in background thread."""
    _pipeline_status["running"] = True
    _pipeline_status["paused"] = False
    _pipeline_status["current_stage"] = ""
    _pipeline_status["last_run"] = datetime.now().isoformat()
    try:
        from src.pipeline import run_scan_only
        result = run_scan_only(sources=sources, date_from=date_from, date_to=date_to)
        _pipeline_status["last_result"] = result
    except Exception as e:
        _pipeline_status["last_result"] = {"status": "failed", "error": str(e)}
    finally:
        _pipeline_status["running"] = False
        _pipeline_status["paused"] = False
        _pipeline_status["current_stage"] = ""


def _run_send_bg(date_from, date_to, bu_filter, industry_filter, min_score,
                 recipient_emails, pdf_mode, preview_only):
    """Run send-only pipeline in background thread."""
    _pipeline_status["running"] = True
    _pipeline_status["paused"] = False
    _pipeline_status["current_stage"] = ""
    _pipeline_status["last_run"] = datetime.now().isoformat()
    try:
        from src.pipeline import run_send_only
        result = run_send_only(
            date_from=date_from, date_to=date_to, bu_filter=bu_filter,
            industry_filter=industry_filter, min_score=min_score,
            recipient_emails=recipient_emails, pdf_mode=pdf_mode,
            preview_only=preview_only,
        )
        _pipeline_status["last_result"] = result
    except Exception as e:
        _pipeline_status["last_result"] = {"status": "failed", "error": str(e)}
    finally:
        _pipeline_status["running"] = False
        _pipeline_status["paused"] = False
        _pipeline_status["current_stage"] = ""


class ScanRequest(BaseModel):
    sources: list[str] | None = None
    date_from: str | None = None
    date_to: str | None = None


@app.post("/api/pipeline/scan")
def run_scan(req: ScanRequest = ScanRequest()):
    """Run scan only (collect + validate + score). No email delivery."""
    if _pipeline_status["running"]:
        raise HTTPException(409, "Pipeline is already running")

    thread = threading.Thread(
        target=_run_scan_bg,
        args=(req.sources, req.date_from, req.date_to),
        daemon=True,
    )
    thread.start()
    return {
        "message": "Scan started",
        "sources": req.sources or ["rss", "reddit", "trends"],
        "status_url": "/api/pipeline/status",
    }


class SendRequest(BaseModel):
    date_from: str | None = None
    date_to: str | None = None
    bu_filter: str | None = None
    industry_filter: str | None = None
    min_score: float = 5.0
    recipient_emails: list[str] | None = None
    pdf_mode: bool = False
    preview_only: bool = False


@app.post("/api/pipeline/send")
def run_send(req: SendRequest = SendRequest()):
    """Compose and send digest from already-scored signals."""
    if _pipeline_status["running"]:
        raise HTTPException(409, "Pipeline is already running")

    thread = threading.Thread(
        target=_run_send_bg,
        args=(req.date_from, req.date_to, req.bu_filter, req.industry_filter,
              req.min_score, req.recipient_emails, req.pdf_mode, req.preview_only),
        daemon=True,
    )
    thread.start()
    return {
        "message": f"{'Preview' if req.preview_only else 'Send'} started",
        "status_url": "/api/pipeline/status",
    }


@app.get("/api/pipeline/scan-log")
def get_scan_log(limit: int = 20):
    """Get scan run history."""
    conn = get_connection()
    try:
        rows = conn.execute(
            "SELECT * FROM scan_log ORDER BY started_at DESC LIMIT ?", (limit,)
        ).fetchall()
        return {"scans": [dict(r) for r in rows]}
    finally:
        conn.close()


# ── Digest History ───────────────────────────────────────────────────

@app.get("/api/digests")
def list_digests(start_date: str | None = None, end_date: str | None = None):
    """List generated digests (HTML and PDF), optionally filtered by date range."""
    MOCK_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    digests = []
    for f in sorted(MOCK_OUTPUT_DIR.glob("digest-*.*"), reverse=True):
        stat = f.stat()
        created = datetime.fromtimestamp(stat.st_mtime)
        # Timeframe filter
        if start_date and created.strftime("%Y-%m-%d") < start_date:
            continue
        if end_date and created.strftime("%Y-%m-%d") > end_date:
            continue
        digests.append({
            "filename": f.name,
            "format": f.suffix.lstrip("."),
            "size_kb": round(stat.st_size / 1024, 1),
            "created_at": created.isoformat(),
            "preview_url": f"/api/digests/{f.name}/preview",
        })
    return {"digests": digests}


@app.get("/api/digests/{filename}/preview")
def preview_digest(filename: str):
    """Preview a generated digest (HTML rendered, PDF downloaded)."""
    path = MOCK_OUTPUT_DIR / filename
    if not path.exists():
        raise HTTPException(404, "Digest not found")

    if path.suffix == ".pdf":
        return FileResponse(
            str(path),
            media_type="application/pdf",
            filename=path.name,
        )

    return HTMLResponse(path.read_text(encoding="utf-8"))


# ── Dashboard Stats ──────────────────────────────────────────────────

@app.get("/api/dashboard")
def dashboard(start_date: str | None = None, end_date: str | None = None):
    """Get dashboard summary stats, optionally filtered by date range."""
    try:
        conn = get_connection()
        init_db()

        # Build date-filtered signal counts
        sig_query = "SELECT COUNT(*) FROM signals"
        sig_scored_query = "SELECT COUNT(*) FROM signals WHERE status='scored'"
        run_query = "SELECT COUNT(*) FROM pipeline_runs"
        params: list = []
        scored_params: list = []
        run_params: list = []

        if start_date or end_date:
            sig_clauses = []
            scored_clauses = ["status='scored'"]
            run_clauses = []
            if start_date:
                sig_clauses.append("collected_at >= ?")
                params.append(start_date)
                scored_clauses.append("collected_at >= ?")
                scored_params.append(start_date)
                run_clauses.append("started_at >= ?")
                run_params.append(start_date)
            if end_date:
                sig_clauses.append("collected_at <= ?")
                params.append(end_date + " 23:59:59")
                scored_clauses.append("collected_at <= ?")
                scored_params.append(end_date + " 23:59:59")
                run_clauses.append("started_at <= ?")
                run_params.append(end_date + " 23:59:59")
            if sig_clauses:
                sig_query += " WHERE " + " AND ".join(sig_clauses)
            sig_scored_query = "SELECT COUNT(*) FROM signals WHERE " + " AND ".join(scored_clauses)
            if run_clauses:
                run_query += " WHERE " + " AND ".join(run_clauses)

        signals_total = conn.execute(sig_query, params).fetchone()[0]
        signals_scored = conn.execute(sig_scored_query, scored_params).fetchone()[0]
        pipeline_runs_count = conn.execute(run_query, run_params).fetchone()[0]
        last_run = conn.execute(
            "SELECT started_at, status FROM pipeline_runs ORDER BY started_at DESC LIMIT 1"
        ).fetchone()

        # Industry and keyword counts (V2.1)
        try:
            industries_count = conn.execute("SELECT COUNT(*) FROM industries WHERE active = 1").fetchone()[0]
            keywords_count = conn.execute("SELECT COUNT(*) FROM keywords WHERE active = 1").fetchone()[0]
        except Exception:
            industries_count = 0
            keywords_count = 0

        conn.close()

        recipients_config = get_recipients()
        active_recipients = sum(
            1 for r in recipients_config.get("recipients", [])
            if r.get("status") == "active"
        )

        digests = list(MOCK_OUTPUT_DIR.glob("digest-*.*")) if MOCK_OUTPUT_DIR.exists() else []

        return {
            "signals_total": signals_total,
            "signals_scored": signals_scored,
            "pipeline_runs": pipeline_runs_count,
            "last_run": {
                "time": last_run[0] if last_run else None,
                "status": last_run[1] if last_run else None,
            },
            "active_recipients": active_recipients,
            "digests_generated": len(digests),
            "pipeline_running": _pipeline_status["running"],
            "industries_count": industries_count,
            "keywords_count": keywords_count,
        }
    except Exception as e:
        logger.error("Dashboard query failed: %s", e)
        return {
            "signals_total": 0,
            "signals_scored": 0,
            "pipeline_runs": 0,
            "last_run": {"time": None, "status": None},
            "active_recipients": 0,
            "digests_generated": 0,
            "pipeline_running": False,
            "industries_count": 0,
            "keywords_count": 0,
        }


@app.get("/api/pipeline/runs")
def list_pipeline_runs(start_date: str | None = None, end_date: str | None = None, limit: int = 50):
    """Get pipeline run history with optional date range filter."""
    conn = get_connection()
    try:
        runs = get_pipeline_runs_by_timeframe(conn, start_date, end_date, limit)
        return {"runs": runs}
    finally:
        conn.close()


@app.get("/api/signals")
def list_signals(start_date: str | None = None, end_date: str | None = None,
                 status: str | None = None, limit: int = 100):
    """Get signals with optional date range and status filters."""
    conn = get_connection()
    try:
        signals = get_signals_by_timeframe(conn, start_date, end_date, status)
        return {"signals": signals[:limit], "total": len(signals)}
    finally:
        conn.close()


# ── Recommendations (V2.1 Phase D) ────────────────────────────────

@app.get("/api/recommendations")
def get_recommendations(max_recommendations: int = 15):
    """Generate AI-powered strategic recommendations from signal patterns."""
    from src.analyzer.recommendations import generate_recommendations
    conn = get_connection()
    try:
        return generate_recommendations(conn, max_recommendations=max_recommendations)
    finally:
        conn.close()


# ── Pattern Detection (V2.1 Phase D) ─────────────────────────────

@app.get("/api/patterns")
def get_patterns():
    """Detect recurring patterns across signals, competitors, and trends."""
    from src.analyzer.pattern_detector import detect_patterns
    conn = get_connection()
    try:
        return detect_patterns(conn)
    finally:
        conn.close()


# ── Admin / Sync (V2.3) ─────────────────────────────────────────────

@app.post("/api/admin/sync-config")
def sync_config():
    """Force-sync industries, keywords, and subreddits from config files into DB.

    Use this if startup sync failed or after editing config files while
    the server is running.
    """
    conn = get_connection()
    try:
        _migrate_v23(conn)
        _auto_seed_industries(conn)
        ind_count = conn.execute("SELECT COUNT(*) FROM industries").fetchone()[0]
        kw_count = conn.execute("SELECT COUNT(*) FROM keywords").fetchone()[0]
        sub_count = conn.execute("SELECT COUNT(*) FROM reddit_subreddits").fetchone()[0]
        return {
            "status": "synced",
            "industries": ind_count,
            "keywords": kw_count,
            "subreddits": sub_count,
        }
    finally:
        conn.close()


# ── Signal Status Management (V2.3) ─────────────────────────────────

@app.post("/api/signals/{signal_id}/dismiss")
def dismiss_signal(signal_id: int):
    """Mark a signal as not-relevant. Removes from views/stats but kept in DB."""
    conn = get_connection()
    try:
        row = conn.execute("SELECT id FROM signals WHERE id = ?", (signal_id,)).fetchone()
        if not row:
            raise HTTPException(404, f"Signal {signal_id} not found")
        conn.execute(
            "UPDATE signals SET dismissed = 1, dismissed_at = datetime('now') WHERE id = ?",
            (signal_id,)
        )
        conn.commit()
        return {"signal_id": signal_id, "dismissed": True}
    finally:
        conn.close()


@app.post("/api/signals/{signal_id}/restore")
def restore_signal(signal_id: int):
    """Restore a previously dismissed signal."""
    conn = get_connection()
    try:
        conn.execute(
            "UPDATE signals SET dismissed = 0, dismissed_at = NULL WHERE id = ?",
            (signal_id,)
        )
        conn.commit()
        return {"signal_id": signal_id, "dismissed": False}
    finally:
        conn.close()


@app.post("/api/signals/{signal_id}/handle")
def handle_signal(signal_id: int, handled_by: str = ""):
    """Mark a signal as read/handled for review tracking."""
    conn = get_connection()
    try:
        row = conn.execute("SELECT id FROM signals WHERE id = ?", (signal_id,)).fetchone()
        if not row:
            raise HTTPException(404, f"Signal {signal_id} not found")
        conn.execute(
            "UPDATE signals SET handled = 1, handled_at = datetime('now'), handled_by = ? WHERE id = ?",
            (handled_by, signal_id)
        )
        conn.commit()
        return {"signal_id": signal_id, "handled": True, "handled_by": handled_by}
    finally:
        conn.close()


@app.post("/api/signals/{signal_id}/unhandle")
def unhandle_signal(signal_id: int):
    """Remove handled mark from a signal."""
    conn = get_connection()
    try:
        conn.execute(
            "UPDATE signals SET handled = 0, handled_at = NULL, handled_by = '' WHERE id = ?",
            (signal_id,)
        )
        conn.commit()
        return {"signal_id": signal_id, "handled": False}
    finally:
        conn.close()


# ── Competitor Pulse (V2.3) ─────────────────────────────────────────

MONITORED_COMPETITORS = [
    "TT Electronics", "HBK", "Hottinger", "Zemic", "Rice Lake",
    "Kistler", "Kyowa", "NMB", "Omega", "Novanta",
    "Flintec", "Sunrise Instruments", "Figure AI", "Boston Dynamics",
    "Mettler Toledo", "Siemens", "Honeywell",
]


@app.get("/api/executive/competitor-pulse")
def competitor_pulse(start_date: str | None = None, end_date: str | None = None,
                     bu_id: str | None = None, industry_id: str | None = None):
    """Top competitors by signal count with trend direction."""
    conn = get_connection()
    try:
        results = []
        for comp in MONITORED_COMPETITORS:
            query = """
                SELECT COUNT(*), AVG(sa.score_composite)
                FROM signals s
                JOIN signal_analysis sa ON s.id = sa.signal_id
                WHERE s.status IN ('scored', 'published')
                  AND COALESCE(s.dismissed, 0) = 0
                  AND (LOWER(s.title) LIKE ? OR LOWER(s.raw_content) LIKE ?
                       OR LOWER(sa.headline) LIKE ? OR LOWER(sa.what_summary) LIKE ?)
            """
            like = f"%{comp.lower()}%"
            params = [like, like, like, like]
            if start_date:
                query += " AND s.collected_at >= ?"
                params.append(start_date)
            if end_date:
                query += " AND s.collected_at <= ?"
                params.append(end_date + " 23:59:59")
            if bu_id:
                query += " AND s.id IN (SELECT signal_id FROM signal_bus WHERE bu_id = ?)"
                params.append(bu_id)
            if industry_id:
                query += " AND s.id IN (SELECT signal_id FROM signal_industries WHERE industry_id = ?)"
                params.append(industry_id)

            row = conn.execute(query, params).fetchone()
            count = row[0] or 0
            if count == 0:
                continue

            # Trend: compare recent 2 weeks vs prior 2 weeks
            recent = conn.execute(
                "SELECT COUNT(*) FROM signals s JOIN signal_analysis sa ON s.id = sa.signal_id "
                "WHERE s.status IN ('scored','published') AND COALESCE(s.dismissed,0)=0 "
                "AND (LOWER(s.title) LIKE ? OR LOWER(sa.headline) LIKE ?) "
                "AND s.collected_at >= date('now', '-14 days')",
                (like, like)
            ).fetchone()[0]
            older = conn.execute(
                "SELECT COUNT(*) FROM signals s JOIN signal_analysis sa ON s.id = sa.signal_id "
                "WHERE s.status IN ('scored','published') AND COALESCE(s.dismissed,0)=0 "
                "AND (LOWER(s.title) LIKE ? OR LOWER(sa.headline) LIKE ?) "
                "AND s.collected_at >= date('now', '-28 days') AND s.collected_at < date('now', '-14 days')",
                (like, like)
            ).fetchone()[0]

            if recent > older:
                trend = "up"
            elif recent < older:
                trend = "down"
            else:
                trend = "stable"

            results.append({
                "competitor": comp,
                "signal_count": count,
                "avg_score": round(row[1] or 0, 1),
                "trend": trend,
                "recent_signals": recent,
            })

        results.sort(key=lambda x: x["signal_count"], reverse=True)
        return {"competitors": results[:10], "total_monitored": len(MONITORED_COMPETITORS)}
    finally:
        conn.close()


# ── Pipeline Last Run (V2.3) ────────────────────────────────────────

@app.get("/api/pipeline/last-run")
def pipeline_last_run():
    """Get the timestamp of the last completed pipeline run."""
    conn = get_connection()
    try:
        row = conn.execute(
            "SELECT started_at, completed_at, status, run_type, signals_collected, signals_scored "
            "FROM pipeline_runs ORDER BY started_at DESC LIMIT 1"
        ).fetchone()
        if not row:
            return {"last_run": None}
        return {
            "last_run": {
                "started_at": row[0], "completed_at": row[1],
                "status": row[2], "run_type": row[3],
                "signals_collected": row[4], "signals_scored": row[5],
            }
        }
    finally:
        conn.close()


# ── Reddit Subreddit Management (V2.3) ─────────────────────────────

class SubredditCreate(BaseModel):
    name: str
    category: str = ""
    notes: str = ""


class SubredditUpdate(BaseModel):
    category: str | None = None
    active: bool | None = None
    notes: str | None = None


@app.get("/api/reddit/subreddits")
def list_subreddits():
    """Get all managed subreddits with signal counts."""
    conn = get_connection()
    try:
        rows = conn.execute(
            "SELECT id, name, category, active, notes, created_at, updated_at "
            "FROM reddit_subreddits ORDER BY category, name"
        ).fetchall()

        # Count valid signals per subreddit (source_name = 'Reddit r/<name>')
        signal_counts = {}
        count_rows = conn.execute(
            "SELECT source_name, COUNT(*) FROM signals "
            "WHERE source_name LIKE 'Reddit r/%' "
            "  AND status IN ('scored', 'published') "
            "  AND COALESCE(dismissed, 0) = 0 "
            "GROUP BY source_name"
        ).fetchall()
        for cr in count_rows:
            # Extract subreddit name from 'Reddit r/name'
            sub_name = cr[0].replace("Reddit r/", "")
            signal_counts[sub_name.lower()] = cr[1]

        subs = []
        for r in rows:
            subs.append({
                "id": r[0], "name": r[1], "category": r[2],
                "active": bool(r[3]), "notes": r[4],
                "created_at": r[5], "updated_at": r[6],
                "signal_count": signal_counts.get(r[1].lower(), 0),
            })
        return {"subreddits": subs, "total": len(subs)}
    finally:
        conn.close()


@app.post("/api/reddit/subreddits")
def add_subreddit(sub: SubredditCreate):
    """Add a new subreddit to monitor."""
    conn = get_connection()
    try:
        existing = conn.execute(
            "SELECT id FROM reddit_subreddits WHERE name = ?", (sub.name,)
        ).fetchone()
        if existing:
            raise HTTPException(400, f"Subreddit r/{sub.name} already exists")
        conn.execute(
            "INSERT INTO reddit_subreddits (name, category, notes) VALUES (?, ?, ?)",
            (sub.name, sub.category, sub.notes)
        )
        conn.commit()
        return {"name": sub.name, "status": "added"}
    finally:
        conn.close()


@app.put("/api/reddit/subreddits/{sub_id}")
def update_subreddit(sub_id: int, updates: SubredditUpdate):
    """Update a subreddit's settings."""
    conn = get_connection()
    try:
        row = conn.execute("SELECT id FROM reddit_subreddits WHERE id = ?", (sub_id,)).fetchone()
        if not row:
            raise HTTPException(404, f"Subreddit {sub_id} not found")
        update_data = updates.model_dump(exclude_none=True)
        if "active" in update_data:
            update_data["active"] = int(update_data["active"])
        for key, val in update_data.items():
            conn.execute(
                f"UPDATE reddit_subreddits SET {key} = ?, updated_at = datetime('now') WHERE id = ?",
                (val, sub_id)
            )
        conn.commit()
        return {"id": sub_id, "updated": True}
    finally:
        conn.close()


@app.delete("/api/reddit/subreddits/{sub_id}")
def delete_subreddit(sub_id: int):
    """Remove a subreddit from monitoring."""
    conn = get_connection()
    try:
        deleted = conn.execute("DELETE FROM reddit_subreddits WHERE id = ?", (sub_id,)).rowcount
        conn.commit()
        if not deleted:
            raise HTTPException(404, f"Subreddit {sub_id} not found")
        return {"deleted": sub_id}
    finally:
        conn.close()


@app.post("/api/reddit/test-collection/{subreddit_name}")
def test_reddit_collection(subreddit_name: str):
    """Test Reddit collection for a single subreddit. Returns detailed results.

    Used by the 'Test Collection' button in the UI for debugging.
    """
    from src.collector.reddit_collector import test_single_subreddit
    result = test_single_subreddit(subreddit_name)
    return result


@app.get("/api/reddit/status")
def reddit_status():
    """Check Reddit API credential and connection status."""
    import os
    status = {
        "praw_installed": False,
        "credentials_configured": False,
        "credential_hint": "",
    }
    try:
        import praw  # noqa: F401
        status["praw_installed"] = True
    except ImportError:
        pass

    client_id = os.getenv("REDDIT_CLIENT_ID", "")
    client_secret = os.getenv("REDDIT_CLIENT_SECRET", "")
    if client_id and client_secret and not client_id.startswith("your-"):
        status["credentials_configured"] = True
        status["credential_hint"] = f"{client_id[:4]}...{client_id[-3:]}"
    else:
        status["credential_hint"] = "Not configured"

    return status


# ── Recommendations Export (V2.3) ───────────────────────────────────

@app.get("/api/export/recommendations/excel")
def export_recommendations_excel():
    """Export recommendations as Excel workbook."""
    try:
        import openpyxl  # noqa: F401
    except ImportError:
        raise HTTPException(422, detail="Install openpyxl: pip install openpyxl")
    from src.analyzer.recommendations import generate_recommendations
    from fastapi.responses import StreamingResponse

    conn = get_connection()
    try:
        recs_data = generate_recommendations(conn)
    finally:
        conn.close()

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Recommendations"
    ws.append(["Priority", "Type", "Title", "Description", "Action", "Owner", "Score"])
    priority_map = {1: "Critical", 2: "High", 3: "Medium"}
    for rec in recs_data.get("recommendations", []):
        ws.append([
            priority_map.get(rec.get("priority", 3), "Medium"),
            rec.get("type", ""), rec.get("title", ""),
            rec.get("description", ""), rec.get("action", ""),
            rec.get("owner", ""), rec.get("score", 0),
        ])

    from openpyxl.styles import Font, PatternFill
    navy_fill = PatternFill(start_color="1B2A4A", end_color="1B2A4A", fill_type="solid")
    for cell in ws[1]:
        cell.fill = navy_fill
        cell.font = Font(name="Arial", bold=True, color="FFFFFF", size=10)
    for col in ws.columns:
        max_len = max(len(str(c.value or "")) for c in col)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 60)

    buffer = BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    filename = f"vpg-recommendations-{datetime.now().strftime('%Y%m%d')}.xlsx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


@app.get("/api/export/recommendations/pptx")
def export_recommendations_pptx():
    """Export recommendations as PowerPoint presentation."""
    try:
        import pptx as pptx_mod  # noqa: F401
    except ImportError:
        raise HTTPException(422, detail="Install python-pptx: pip install python-pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from src.analyzer.recommendations import generate_recommendations
    from fastapi.responses import StreamingResponse

    conn = get_connection()
    try:
        recs_data = generate_recommendations(conn)
    finally:
        conn.close()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(27, 42, 74)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "VPG Strategic Recommendations"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p2 = tb.text_frame.add_paragraph()
    p2.text = f"Generated {datetime.now().strftime('%B %d, %Y')}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(46, 117, 182)

    priority_colors = {1: RGBColor(229, 57, 53), 2: RGBColor(245, 124, 0), 3: RGBColor(46, 117, 182)}
    priority_labels = {1: "CRITICAL", 2: "HIGH", 3: "MEDIUM"}

    for rec in recs_data.get("recommendations", [])[:12]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        pri = rec.get("priority", 3)
        # Priority badge
        hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        hp = hdr.text_frame.paragraphs[0]
        r1 = hp.add_run()
        r1.text = f"{priority_labels.get(pri, 'MEDIUM')}  "
        r1.font.size = Pt(11)
        r1.font.bold = True
        r1.font.color.rgb = priority_colors.get(pri, RGBColor(46, 117, 182))
        r2 = hp.add_run()
        r2.text = rec.get("type", "")
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(100, 100, 100)

        # Title
        ttl = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
        ttl.text_frame.word_wrap = True
        tp = ttl.text_frame.paragraphs[0]
        tp.text = rec.get("title", "")
        tp.font.size = Pt(20)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(27, 42, 74)

        # Description
        desc = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
        desc.text_frame.word_wrap = True
        dp = desc.text_frame.paragraphs[0]
        dp.text = rec.get("description", "")[:400]
        dp.font.size = Pt(12)
        dp.font.color.rgb = RGBColor(60, 60, 60)

        # Action
        act = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
        act.text_frame.word_wrap = True
        ap = act.text_frame.paragraphs[0]
        ar1 = ap.add_run()
        ar1.text = "ACTION: "
        ar1.font.size = Pt(12)
        ar1.font.bold = True
        ar1.font.color.rgb = RGBColor(232, 121, 47)
        ar2 = ap.add_run()
        ar2.text = rec.get("action", "")[:300]
        ar2.font.size = Pt(12)
        ar2.font.color.rgb = RGBColor(60, 60, 60)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    filename = f"vpg-recommendations-{datetime.now().strftime('%Y%m%d')}.pptx"
    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )


# ── Phase 3: Feedback & Scoring Refinement ─────────────────────────

@app.get("/api/feedback/summary")
def feedback_summary():
    """Get feedback summary and scoring adjustment data."""
    from src.feedback.scoring_refinement import get_feedback_summary, compute_scoring_adjustments
    conn = get_connection()
    try:
        summary = get_feedback_summary(conn)
        adjustments = compute_scoring_adjustments(conn)
        return {"feedback": summary, "adjustments": adjustments}
    finally:
        conn.close()


@app.post("/api/feedback")
def submit_feedback(body: dict):
    """Submit thumbs-up/down feedback for a signal."""
    signal_id = body.get("signal_id")
    rating = body.get("rating")  # "up" or "down"
    email = body.get("email", "anonymous")
    comment = body.get("comment", "")

    if not signal_id or rating not in ("up", "down"):
        raise HTTPException(400, "signal_id and rating ('up'/'down') required")

    conn = get_connection()
    try:
        conn.execute(
            "INSERT INTO feedback (signal_id, recipient_email, rating, comment) VALUES (?, ?, ?, ?)",
            (signal_id, email, rating, comment),
        )
        conn.commit()
        return {"status": "ok", "signal_id": signal_id, "rating": rating}
    finally:
        conn.close()


# ── Phase 3: Pre-Event Intelligence Packs ──────────────────────────

@app.get("/api/events")
def list_events():
    """List all configured events."""
    from src.events.intel_packs import list_events as _list, get_upcoming_events
    return {"events": _list(), "upcoming": get_upcoming_events(90)}


@app.post("/api/events")
def create_event(body: dict):
    """Add a new event."""
    from src.events.intel_packs import add_event
    required = ["id", "name", "start_date"]
    for field in required:
        if field not in body:
            raise HTTPException(400, f"Missing required field: {field}")
    return add_event(body)


@app.put("/api/events/{event_id}")
def update_event_endpoint(event_id: str, body: dict):
    """Update an existing event."""
    from src.events.intel_packs import update_event
    result = update_event(event_id, body)
    if result is None:
        raise HTTPException(404, f"Event '{event_id}' not found")
    return result


@app.delete("/api/events/{event_id}")
def delete_event_endpoint(event_id: str):
    """Delete an event."""
    from src.events.intel_packs import delete_event
    if not delete_event(event_id):
        raise HTTPException(404, f"Event '{event_id}' not found")
    return {"status": "deleted", "event_id": event_id}


@app.get("/api/events/{event_id}/intel-pack")
def get_intel_pack(event_id: str):
    """Generate a pre-event intelligence pack."""
    from src.events.intel_packs import generate_intel_pack
    conn = get_connection()
    try:
        return generate_intel_pack(event_id, conn)
    finally:
        conn.close()


# ── Phase 3: Outreach Templates ────────────────────────────────────

@app.get("/api/outreach/{signal_id}")
def get_outreach_templates(signal_id: int):
    """Generate outreach templates for a signal."""
    from src.outreach.templates import generate_outreach
    conn = get_connection()
    try:
        return generate_outreach(signal_id, conn)
    finally:
        conn.close()


@app.post("/api/outreach/batch")
def batch_outreach(body: dict):
    """Generate outreach templates for multiple signals."""
    from src.outreach.templates import generate_batch_outreach
    signal_ids = body.get("signal_ids", [])
    if not signal_ids:
        raise HTTPException(400, "signal_ids list required")
    conn = get_connection()
    try:
        return {"results": generate_batch_outreach(signal_ids, conn)}
    finally:
        conn.close()


# ── Phase 3: India Production Advantage Monitor ────────────────────

@app.get("/api/india/monitor")
def india_monitor():
    """Get India production advantage intelligence."""
    from src.india.monitor import analyze_india_signals
    conn = get_connection()
    try:
        return analyze_india_signals(conn)
    finally:
        conn.close()


@app.get("/api/india/talking-points/{signal_id}")
def india_talking_points(signal_id: int):
    """Get India-specific talking points for a signal."""
    from src.india.monitor import get_india_talking_points_for_signal
    conn = get_connection()
    try:
        return get_india_talking_points_for_signal(signal_id, conn)
    finally:
        conn.close()


# ── Phase 3: Monthly Effectiveness Reports ─────────────────────────

@app.get("/api/reports/monthly")
def monthly_report(year: int = None, month: int = None):
    """Generate a monthly effectiveness report."""
    from src.reports.monthly import generate_monthly_report
    conn = get_connection()
    try:
        return generate_monthly_report(year, month, conn)
    finally:
        conn.close()


# ── Phase 3: Meeting Prep Briefs ───────────────────────────────────

@app.get("/api/accounts")
def list_target_accounts():
    """List all configured target accounts."""
    from src.reports.meeting_prep import list_target_accounts
    return {"accounts": list_target_accounts()}


@app.get("/api/accounts/{account_key}/meeting-brief")
def get_meeting_brief(account_key: str):
    """Generate a meeting prep brief for a target account."""
    from src.reports.meeting_prep import generate_meeting_brief
    conn = get_connection()
    try:
        result = generate_meeting_brief(account_key, conn)
        if "error" in result:
            raise HTTPException(404, result["error"])
        return result
    finally:
        conn.close()


# ── Serve React Frontend ────────────────────────────────────────────

UI_BUILD_DIR = PROJECT_ROOT / "src" / "ui" / "build"

if UI_BUILD_DIR.exists():
    # Vite builds to "assets/", Create React App builds to "static/"
    for static_dir_name in ("assets", "static"):
        static_dir = UI_BUILD_DIR / static_dir_name
        if static_dir.exists():
            app.mount(f"/{static_dir_name}", StaticFiles(directory=str(static_dir)), name=static_dir_name)

    @app.get("/{full_path:path}")
    def serve_react(full_path: str):
        """Serve the React SPA for any non-API route."""
        file_path = UI_BUILD_DIR / full_path
        if file_path.exists() and file_path.is_file():
            return FileResponse(str(file_path))
        return FileResponse(str(UI_BUILD_DIR / "index.html"))


def start_server(host: str = "0.0.0.0", port: int = 8000):
    """Start the API server."""
    import uvicorn
    print(f"\n  VPG Intelligence Digest — Management UI")
    print(f"  API:   http://localhost:{port}/api/dashboard")
    print(f"  UI:    http://localhost:{port}")
    print(f"  Docs:  http://localhost:{port}/docs\n")
    uvicorn.run(app, host=host, port=port)


if __name__ == "__main__":
    start_server()
